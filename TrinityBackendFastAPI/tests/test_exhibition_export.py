import base64
import importlib.util
import io
import sys
import types
from pathlib import Path

import pytest
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


BACKEND_ROOT = Path(__file__).resolve().parents[1]
EXHIBITION_PATH = BACKEND_ROOT / "app" / "features" / "exhibition"


def _load_module(module_name: str, file_path: Path):
    spec = importlib.util.spec_from_file_location(module_name, file_path)
    module = importlib.util.module_from_spec(spec)
    assert spec and spec.loader  # pragma: no cover - defensive
    sys.modules[module_name] = module
    spec.loader.exec_module(module)  # type: ignore[attr-defined]
    return module


# Create lightweight package hierarchy to satisfy relative imports without
# executing app.main side effects.
app_pkg = types.ModuleType("app")
app_pkg.__path__ = [str(BACKEND_ROOT / "app")]
features_pkg = types.ModuleType("app.features")
features_pkg.__path__ = [str(BACKEND_ROOT / "app" / "features")]
exhibition_pkg = types.ModuleType("app.features.exhibition")
exhibition_pkg.__path__ = [str(EXHIBITION_PATH)]

sys.modules.setdefault("app", app_pkg)
sys.modules.setdefault("app.features", features_pkg)
sys.modules.setdefault("app.features.exhibition", exhibition_pkg)

schemas_module = _load_module("app.features.exhibition.schemas", EXHIBITION_PATH / "schemas.py")
export_module = _load_module("app.features.exhibition.export", EXHIBITION_PATH / "export.py")

DocumentStylesPayload = schemas_module.DocumentStylesPayload
ExhibitionExportRequest = schemas_module.ExhibitionExportRequest
ExportGenerationError = export_module.ExportGenerationError
SlideDomSnapshotPayload = schemas_module.SlideDomSnapshotPayload
build_export_filename = export_module.build_export_filename
build_pdf_bytes = export_module.build_pdf_bytes
build_pptx_bytes = export_module.build_pptx_bytes
_px_to_emu = export_module._px_to_emu


# 1x1 transparent PNG
def _png_data_url() -> str:
    buffer = io.BytesIO()
    Image.new("RGB", (10, 10), color="#3366FF").save(buffer, format="PNG")
    pixel = base64.b64encode(buffer.getvalue()).decode("ascii")
    return f"data:image/png;base64,{pixel}"


def _build_payload(include_screenshot: bool = True):
    screenshot = None
    if include_screenshot:
        screenshot = {
            "dataUrl": _png_data_url(),
            "width": 960,
            "height": 540,
            "cssWidth": 960,
            "cssHeight": 540,
            "pixelRatio": 2,
        }

    slides = [
        {
            "id": "slide-001",
            "index": 0,
            "title": "Overview",
            "baseWidth": 960,
            "baseHeight": 540,
            "objects": [
                {
                    "id": "text-001",
                    "type": "text-box",
                    "x": 120,
                    "y": 140,
                    "width": 400,
                    "height": 160,
                    "props": {
                        "text": "<strong>Hello</strong> world!",
                        "fontFamily": "Arial",
                        "fontSize": 24,
                        "align": "center",
                    },
                },
                {
                    "id": "image-001",
                    "type": "image",
                    "x": 560,
                    "y": 200,
                    "width": 120,
                    "height": 120,
                    "props": {"src": _png_data_url()},
                },
                {
                    "id": "shape-001",
                    "type": "shape",
                    "x": 320,
                    "y": 120,
                    "width": 180,
                    "height": 180,
                    "props": {
                        "shapeId": "rounded-rectangle",
                        "fill": "#7C3AED",
                        "stroke": "#312E81",
                        "strokeWidth": 8,
                        "opacity": 0.9,
                    },
                },
            ],
            "screenshot": screenshot,
        },
        {
            "id": "slide-002",
            "index": 1,
            "title": "Fallback Dimensions",
            "baseWidth": 0,
            "baseHeight": 0,
            "objects": [],
            "screenshot": screenshot,
        },
    ]

    return ExhibitionExportRequest.model_validate({"title": "Demo", "slides": slides})


def test_build_pptx_bytes_renders_slides(tmp_path: Path) -> None:
    payload = _build_payload()

    pptx_bytes = build_pptx_bytes(payload)
    assert pptx_bytes.startswith(b"PK")

    from pptx import Presentation

    presentation = Presentation(io.BytesIO(pptx_bytes))
    assert len(presentation.slides) == 2

    text_shapes = [
        shape
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text") and "Hello" in shape.text
    ]
    assert text_shapes, "Expected text box to contain exported content"


def test_build_pptx_bytes_includes_shape() -> None:
    payload = _build_payload()

    pptx_bytes = build_pptx_bytes(payload)
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(pptx_bytes))
    shape_fills = [
        getattr(getattr(shape.fill, "fore_color", None), "rgb", None)
        for shape in presentation.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]

    assert any(fill == RGBColor(0x7C, 0x3A, 0xED) for fill in shape_fills)


def test_build_pptx_bytes_uses_max_dimensions() -> None:
    payload = _build_payload()

    payload.slides[0].base_width = 840
    payload.slides[0].base_height = 520
    payload.slides[1].base_width = 1088
    payload.slides[1].base_height = 520

    pptx_bytes = build_pptx_bytes(payload)
    assert pptx_bytes.startswith(b"PK")

    from pptx import Presentation

    presentation = Presentation(io.BytesIO(pptx_bytes))
    assert presentation.slide_width == _px_to_emu(1088)
    assert presentation.slide_height == _px_to_emu(520)


def test_build_pdf_bytes_requires_screenshots() -> None:
    payload = _build_payload(include_screenshot=True)
    pdf_bytes = build_pdf_bytes(payload)
    assert pdf_bytes.startswith(b"%PDF")

    payload_missing = _build_payload(include_screenshot=False)
    with pytest.raises(ExportGenerationError):
        build_pdf_bytes(payload_missing)


def test_build_pdf_bytes_renders_missing_screenshots(monkeypatch) -> None:
    payload = _build_payload(include_screenshot=False)

    for slide in payload.slides:
        slide.dom_snapshot = SlideDomSnapshotPayload.model_validate(
            {
                "html": "<div class=\"slide\">Slide</div>",
                "width": 960,
                "height": 540,
                "pixelRatio": 2,
            }
        )

    payload.document_styles = DocumentStylesPayload.model_validate(
        {
            "inline": [".slide { width: 100%; height: 100%; }"],
            "external": [],
            "baseUrl": "http://localhost:3000",
        }
    )

    captured_calls: list[tuple[bool, list[str]]] = []

    def fake_request(slides, styles, *, strict=True):
        assert isinstance(styles.inline, list)
        captured_calls.append((strict, [slide.id for slide in slides]))
        if not strict:
            return {}
        return {
            slide.id: {
                "id": slide.id,
                "dataUrl": _png_data_url(),
                "width": 960,
                "height": 540,
                "cssWidth": 960,
                "cssHeight": 540,
                "pixelRatio": 2,
            }
            for slide in slides
        }

    monkeypatch.setattr(export_module, "_request_slide_screenshots", fake_request)

    pdf_bytes = build_pdf_bytes(payload)
    assert pdf_bytes.startswith(b"%PDF")
    assert captured_calls == [
        (False, [slide.id for slide in payload.slides]),
        (True, [slide.id for slide in payload.slides]),
    ]


def test_build_pdf_bytes_falls_back_to_client_capture(monkeypatch) -> None:
    payload = _build_payload(include_screenshot=True)

    payload.document_styles = DocumentStylesPayload.model_validate(
        {
            "inline": [],
            "external": [],
            "baseUrl": "http://localhost:3000",
        }
    )

    class FakeError(ExportGenerationError):
        pass

    def failing_request(slides, styles, *, strict=True):
        raise FakeError("renderer offline")

    monkeypatch.setattr(export_module, "_request_slide_screenshots", failing_request)

    pdf_bytes = build_pdf_bytes(payload)
    assert pdf_bytes.startswith(b"%PDF")


def test_build_export_filename_normalises_title() -> None:
    name = build_export_filename(" Quarterly Results 2024 / Demo  ", "pptx")
    assert name == "quarterly-results-2024-demo.pptx"

    fallback = build_export_filename(None, "pdf")
    assert fallback == "exhibition-export.pdf"
