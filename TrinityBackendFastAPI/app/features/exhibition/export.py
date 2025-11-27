from __future__ import annotations

import base64
import html
import io
import json
import logging
import math
import re
import urllib.parse
from typing import Any, Iterable, Optional, Sequence, Tuple
from urllib.request import urlopen, Request
from urllib.error import URLError, HTTPError
from urllib.parse import urlparse, parse_qs

try:
    from app.DataStorageRetrieval.minio_utils import get_client, MINIO_BUCKET
    from minio.error import S3Error
    MINIO_AVAILABLE = True
except ImportError:
    MINIO_AVAILABLE = False

from pptx import Presentation
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas
from PIL import Image

# Helper function to convert PNG screenshot to JPG if needed
def _convert_png_to_jpg(png_bytes: bytes, quality: int = 95) -> bytes:
    """Convert PNG screenshot bytes to JPG format.
    
    This ensures the same screenshot is used for both JPG and PDF exports.
    The PNG is captured once, then converted to JPG when needed.
    
    Args:
        png_bytes: PNG image bytes from screenshot
        quality: JPG quality (1-100, default 95)
    
    Returns:
        JPG image bytes
    """
    try:
        img = Image.open(io.BytesIO(png_bytes))
        # Convert RGBA to RGB if needed (JPG doesn't support transparency)
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        output = io.BytesIO()
        img.save(output, format='JPEG', quality=quality, optimize=True)
        return output.getvalue()
    except Exception as exc:
        logger.warning('Failed to convert PNG to JPG, using PNG instead: %s', exc)
        return png_bytes  # Fallback to PNG if conversion fails

try:  # pragma: no cover - optional dependency for SVG rasterisation
    import cairosvg  # type: ignore[import-not-found]
except Exception:  # pragma: no cover - optional dependency missing or misconfigured
    cairosvg = None  # type: ignore[assignment]

from .renderer import ExhibitionRendererError, build_inputs, render_slide_batch
from .schemas import (
    DocumentStylesPayload,
    ExhibitionExportRequest,
    SlideDomSnapshotPayload,
    SlideExportObjectPayload,
    SlideExportPayload,
    SlideScreenshotPayload,
    SlideScreenshotResponse,
)

logger = logging.getLogger(__name__)

PX_PER_INCH = 96.0
EMU_PER_INCH = 914400
PT_PER_INCH = 72.0

METADATA_MARKER = "TRINITY_EXPORT_METADATA"

CANVAS_STAGE_HEIGHT = 520.0
TOP_LAYOUT_MIN_HEIGHT = 210.0
BOTTOM_LAYOUT_MIN_HEIGHT = 220.0
SIDE_LAYOUT_MIN_WIDTH = 280.0
SIDE_LAYOUT_RATIO = 0.34

GRADIENT_PRESETS: dict[str, dict[str, object]] = {
    "default": {"stops": ["#7c3aed", "#ec4899", "#f97316"], "angle": 135},
    "blue": {"stops": ["#1d4ed8", "#2563eb", "#0ea5e9", "#14b8a6"], "angle": 135},
    "purple": {"stops": ["#5b21b6", "#7c3aed", "#a855f7", "#ec4899"], "angle": 135},
    "green": {"stops": ["#047857", "#10b981", "#22c55e", "#bef264"], "angle": 135},
    "orange": {"stops": ["#c2410c", "#ea580c", "#f97316", "#facc15"], "angle": 135},
    "gradient-aurora": {"stops": ["#312e81", "#7c3aed", "#ec4899", "#f97316"], "angle": 135},
    "gradient-dusk": {"stops": ["#1e3a8a", "#6366f1", "#a855f7", "#f472b6"], "angle": 135},
    "gradient-oceanic": {"stops": ["#0f172a", "#1d4ed8", "#38bdf8", "#2dd4bf"], "angle": 135},
    "gradient-forest": {"stops": ["#064e3b", "#047857", "#22c55e", "#a3e635"], "angle": 135},
    "gradient-tropical": {"stops": ["#0ea5e9", "#22d3ee", "#34d399", "#fde68a"], "angle": 135},
    "gradient-blush": {"stops": ["#f472b6", "#fb7185", "#f97316", "#fde68a"], "angle": 135},
    "gradient-midnight": {"stops": ["#0f172a", "#312e81", "#6d28d9", "#a855f7"], "angle": 135},
}

DEFAULT_OVERLAY_COLOR = "#7c3aed"
DEFAULT_ATOM_BACKGROUND = RGBColor(0xF1, 0xF5, 0xF9)
DEFAULT_ATOM_BORDER = RGBColor(0xD1, 0xD5, 0xDB)

SHAPE_ID_TO_MSO: dict[str, MSO_SHAPE] = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded-rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "star": MSO_SHAPE.STAR_5_POINT,
    "burst": MSO_SHAPE.EXPLOSION1,
    "arrow-right": MSO_SHAPE.RIGHT_ARROW,
    "arrow-left": MSO_SHAPE.LEFT_ARROW,
    "arrow-up": MSO_SHAPE.UP_ARROW,
    "arrow-down": MSO_SHAPE.DOWN_ARROW,
    "process": MSO_SHAPE.FLOWCHART_PROCESS,
    "decision": MSO_SHAPE.FLOWCHART_DECISION,
    "terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "data": MSO_SHAPE.FLOWCHART_DATA,
    "speech-rectangle": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "speech-oval": MSO_SHAPE.OVAL_CALLOUT,
    "thought-bubble": MSO_SHAPE.CLOUD_CALLOUT,
    "cloud": MSO_SHAPE.CLOUD,
    "double-cloud": MSO_SHAPE.CLOUD,
}

LINE_SHAPES = {"line-horizontal", "line-vertical", "line-diagonal"}
class ExportGenerationError(Exception):
    """Raised when an export file cannot be generated."""


def _px_to_inches(value: float) -> float:
    return value / PX_PER_INCH


def _px_to_emu(value: float) -> Emu:
    return Emu(int(round(_px_to_inches(value) * EMU_PER_INCH)))


def _px_to_pt(value: float) -> float:
    return _px_to_inches(value) * PT_PER_INCH


def _safe_float(value: Optional[float], default: float = 0.0) -> float:
    if value is None:
        return default
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return default
    if math.isnan(numeric) or math.isinf(numeric):
        return default
    return numeric


def _decode_data_url(data_url: str) -> bytes:
    """Decode a data URL (base64 encoded image).
    
    Args:
        data_url: Data URL string (e.g., "data:image/png;base64,...")
    
    Returns:
        Decoded image bytes
    
    Raises:
        ExportGenerationError: If decoding fails
    """
    if not data_url:
        raise ExportGenerationError('Missing image data for slide screenshot.')

    match = re.match(r"^data:.*?;base64,(.+)$", data_url, flags=re.IGNORECASE | re.DOTALL)
    payload = match.group(1) if match else data_url

    try:
        return base64.b64decode(payload, validate=True)
    except (base64.binascii.Error, ValueError) as exc:  # type: ignore[attr-defined]
        raise ExportGenerationError('Unable to decode base64 image data.') from exc


def _load_image_asset(
    source: str,
    obj_id: Optional[str] = None,
    max_size_mb: float = 10.0
) -> Tuple[bytes, Optional[str]]:
    """
    Load and normalize an image asset from various sources.
    
    This helper function handles:
    - Data URLs (base64 encoded images)
    - Remote URLs (HTTP/HTTPS)
    - Converts unsupported formats (WebP, SVG) to PNG/JPEG
    - Normalizes image formats for PDF/PPTX compatibility
    
    Args:
        source: Image source - can be a data URL, HTTP/HTTPS URL, or base64 string
        obj_id: Optional object ID for logging purposes
        max_size_mb: Maximum image size in MB (default 10MB)
    
    Returns:
        Tuple of (image_bytes, format_hint) where format_hint is 'PNG' or 'JPEG'
    
    Raises:
        ExportGenerationError: If image cannot be loaded or converted
    """
    if not source or not isinstance(source, str):
        raise ExportGenerationError(f'Invalid image source for object {obj_id}: {source is None and "None" or type(source).__name__}')
    
    obj_ref = f'object {obj_id}' if obj_id else 'image'
    logger.debug('Loading image asset for %s from source type: %s', obj_ref, 'URL' if source.startswith(('http://', 'https://')) else 'data URL' if source.startswith('data:') else 'unknown')
    
    image_bytes: bytes
    source_type: str
    
    # Handle data URLs (base64 encoded)
    if source.startswith('data:'):
        try:
            image_bytes = _decode_data_url(source)
            source_type = 'data URL'
            logger.debug('Loaded image from data URL for %s: %d bytes', obj_ref, len(image_bytes))
        except Exception as exc:
            raise ExportGenerationError(f'Failed to decode data URL for {obj_ref}: {exc}') from exc
    
    # Handle remote URLs
    elif source.startswith(('http://', 'https://')):
        # Parse URL once for all checks
        parsed_url = urlparse(source)
        
        # Check if this is an internal image API endpoint - fetch directly from MinIO to bypass auth
        is_internal_image_api = (
            '/api/images/content' in parsed_url.path or 
            '/images/content' in parsed_url.path
        )
        
        if is_internal_image_api and MINIO_AVAILABLE:
            # Extract object_name from query parameter
            query_params = parse_qs(parsed_url.query)
            object_name = query_params.get('object_name', [None])[0]
            
            if object_name:
                try:
                    logger.debug('Fetching image from MinIO for %s: object_name=%s', obj_ref, object_name)
                    client = get_client()
                    response = client.get_object(MINIO_BUCKET, object_name)
                    try:
                        image_bytes = response.read()
                        # Verify actual size
                        actual_size_mb = len(image_bytes) / (1024 * 1024)
                        if actual_size_mb > max_size_mb:
                            raise ExportGenerationError(
                                f'Image from MinIO for {obj_ref} is too large: {actual_size_mb:.2f}MB (max: {max_size_mb}MB)'
                            )
                        source_type = 'MinIO'
                        logger.info('Loaded image from MinIO for %s: %d bytes (%.2f MB)', obj_ref, len(image_bytes), actual_size_mb)
                    finally:
                        response.close()
                        response.release_conn()
                except S3Error as exc:
                    if exc.code == 'NoSuchKey':
                        raise ExportGenerationError(f'Image not found in MinIO for {obj_ref}: {object_name}') from exc
                    raise ExportGenerationError(f'MinIO error fetching image for {obj_ref}: {exc}') from exc
                except Exception as exc:
                    raise ExportGenerationError(f'Unexpected error fetching image from MinIO for {obj_ref}: {exc}') from exc
            else:
                # No object_name in query, fall back to HTTP
                logger.warning('Internal image API URL missing object_name parameter, falling back to HTTP: %s', source[:100])
                is_internal_image_api = False
        
        # Fall back to HTTP fetch for external URLs or if MinIO fetch failed
        if not is_internal_image_api or not MINIO_AVAILABLE:
            try:
                logger.debug('Fetching image from URL for %s: %s', obj_ref, source[:100])
                
                # Create request with proper headers for external image APIs (Pixabay, etc.)
                # Some CDNs require User-Agent and Accept headers
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                    'Accept': 'image/webp,image/apng,image/*,*/*;q=0.8',
                    'Accept-Language': 'en-US,en;q=0.9',
                    'Referer': parsed_url.scheme + '://' + parsed_url.netloc + '/',
                }
                
                request = Request(source, headers=headers)
                with urlopen(request, timeout=30) as response:
                    # Check content size
                    content_length = response.headers.get('Content-Length')
                    if content_length:
                        size_mb = int(content_length) / (1024 * 1024)
                        if size_mb > max_size_mb:
                            raise ExportGenerationError(
                                f'Image from URL for {obj_ref} is too large: {size_mb:.2f}MB (max: {max_size_mb}MB)'
                            )
                    
                    image_bytes = response.read()
                    # Verify actual size
                    actual_size_mb = len(image_bytes) / (1024 * 1024)
                    if actual_size_mb > max_size_mb:
                        raise ExportGenerationError(
                            f'Image from URL for {obj_ref} is too large: {actual_size_mb:.2f}MB (max: {max_size_mb}MB)'
                        )
                
                source_type = 'URL'
                logger.info('Loaded image from URL for %s: %d bytes (%.2f MB)', obj_ref, len(image_bytes), actual_size_mb)
            except HTTPError as exc:
                raise ExportGenerationError(f'HTTP error fetching image from URL for {obj_ref}: {exc.code} {exc.reason}') from exc
            except URLError as exc:
                raise ExportGenerationError(f'Network error fetching image from URL for {obj_ref}: {exc.reason}') from exc
            except Exception as exc:
                raise ExportGenerationError(f'Unexpected error fetching image from URL for {obj_ref}: {exc}') from exc
    
    # Handle plain base64 strings (fallback)
    else:
        try:
            image_bytes = base64.b64decode(source, validate=True)
            source_type = 'base64 string'
            logger.debug('Loaded image from base64 string for %s: %d bytes', obj_ref, len(image_bytes))
        except Exception as exc:
            raise ExportGenerationError(f'Failed to decode base64 string for {obj_ref}: {exc}') from exc
    
    if not image_bytes:
        raise ExportGenerationError(f'Empty image data for {obj_ref}')
    
    # Normalize image format using PIL
    try:
        img = Image.open(io.BytesIO(image_bytes))
        original_format = img.format
        original_mode = img.mode
        original_size = img.size
        
        logger.debug(
            'Image for %s: format=%s, mode=%s, size=%dx%d',
            obj_ref, original_format, original_mode, original_size[0], original_size[1]
        )
        
        # Convert unsupported formats
        # WebP, SVG, and other formats need conversion for PDF/PPTX compatibility
        needs_conversion = False
        target_format = 'PNG'  # Default to PNG for transparency support
        
        if original_format in ('WEBP', 'SVG', 'BMP', 'TIFF', 'ICO'):
            needs_conversion = True
            logger.info('Converting image for %s from %s to PNG/JPEG for PDF/PPTX compatibility', obj_ref, original_format)
        
        # Handle SVG specially - convert to PNG via cairosvg if available
        if original_format == 'SVG':
            if cairosvg is None:
                raise ExportGenerationError(
                    f'SVG images require cairosvg library for {obj_ref}. '
                    'Please install it: pip install cairosvg'
                )
            try:
                # Convert SVG to PNG
                png_bytes = cairosvg.svg2png(bytestring=image_bytes)
                img = Image.open(io.BytesIO(png_bytes))
                original_format = 'PNG'
                image_bytes = png_bytes
                logger.info('Converted SVG to PNG for %s: %d bytes', obj_ref, len(png_bytes))
            except Exception as exc:
                raise ExportGenerationError(f'Failed to convert SVG to PNG for {obj_ref}: {exc}') from exc
        
        # Convert RGBA/LA/P modes to RGB for JPEG compatibility
        # Keep PNG format for transparency if original had alpha channel
        if img.mode in ('RGBA', 'LA', 'P'):
            if needs_conversion or original_format not in ('PNG', None):
                # Convert to RGB for JPEG
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img)
                img = background
                target_format = 'JPEG'
            # Otherwise keep PNG for transparency
        elif img.mode != 'RGB':
            img = img.convert('RGB')
            target_format = 'JPEG'
        
        # Save to bytes in target format
        if needs_conversion or original_format != target_format:
            output = io.BytesIO()
            if target_format == 'JPEG':
                img.save(output, format='JPEG', quality=95, optimize=True)
            else:
                img.save(output, format='PNG', optimize=True)
            image_bytes = output.getvalue()
            logger.info(
                'Normalized image for %s: %s -> %s, %d bytes',
                obj_ref, original_format or 'unknown', target_format, len(image_bytes)
            )
        else:
            logger.debug('Image for %s already in compatible format: %s', obj_ref, original_format)
        
        return image_bytes, target_format
        
    except Exception as exc:
        # If PIL fails, log but try to use original bytes (might work for some formats)
        logger.warning(
            'Failed to normalize image for %s using PIL (format detection/conversion): %s. '
            'Attempting to use original bytes.',
            obj_ref, exc
        )
        # Return original bytes - let PDF/PPTX libraries handle it
        return image_bytes, None


def _decode_svg_data_url(data_url: str) -> bytes:
    if not data_url:
        raise ExportGenerationError('Missing SVG data for chart render.')

    if ';base64,' in data_url:
        try:
            _, payload = data_url.split(';base64,', 1)
        except ValueError as exc:  # pragma: no cover - defensive
            raise ExportGenerationError('Invalid SVG data URL.') from exc
        try:
            return base64.b64decode(payload, validate=True)
        except (base64.binascii.Error, ValueError) as exc:  # type: ignore[attr-defined]
            raise ExportGenerationError('Unable to decode base64 SVG data.') from exc

    if ',' not in data_url:
        raise ExportGenerationError('Invalid SVG data URL.')

    _, payload = data_url.split(',', 1)
    try:
        return urllib.parse.unquote_to_bytes(payload)
    except Exception as exc:  # pragma: no cover - unexpected percent decoding errors
        raise ExportGenerationError('Unable to decode SVG data URL.') from exc


_BR_TAG_RE = re.compile(r"<\s*br\s*/?\s*>", flags=re.IGNORECASE)
_TAG_RE = re.compile(r"<[^>]+>")


def _html_to_plain_text(raw: str) -> str:
    if not raw:
        return ''
    text = _BR_TAG_RE.sub('\n', raw)
    text = _TAG_RE.sub('', text)
    return html.unescape(text).strip()


def _parse_hex_color(value: Optional[str]) -> Optional[RGBColor]:
    if not value:
        return None
    cleaned = value.strip().lstrip('#')
    if len(cleaned) not in {6, 3}:
        return None
    if len(cleaned) == 3:
        cleaned = ''.join(ch * 2 for ch in cleaned)
    try:
        red = int(cleaned[0:2], 16)
        green = int(cleaned[2:4], 16)
        blue = int(cleaned[4:6], 16)
    except ValueError:
        return None
    return RGBColor(red, green, blue)


def _parse_color_token(value: Optional[str]) -> Optional[RGBColor]:
    if value is None:
        return None

    token = str(value).strip()
    if not token or token.lower() in {"transparent", "none", "currentcolor"}:
        return None

    if token.startswith("var("):
        return None

    return _parse_hex_color(token)


def _lighten_color(color: RGBColor, ratio: float = 0.2) -> RGBColor:
    ratio = max(0.0, min(1.0, ratio))
    red = min(255, int(color[0]) + int((255 - int(color[0])) * ratio))
    green = min(255, int(color[1]) + int((255 - int(color[1])) * ratio))
    blue = min(255, int(color[2]) + int((255 - int(color[2])) * ratio))
    return RGBColor(red, green, blue)


def _is_non_empty_str(value: Any) -> bool:
    return isinstance(value, str) and value.strip() != ""


def _as_dict(value: Any) -> Optional[dict[str, Any]]:
    return value if isinstance(value, dict) else None


def _ensure_list(value: Any) -> list[Any]:
    if isinstance(value, list):
        return value
    return []


def _is_data_url(value: str) -> bool:
    return bool(re.match(r"^data:[^;]+;base64,", value, flags=re.IGNORECASE))


def _humanise_key(value: str) -> str:
    if not value:
        return value
    cleaned = re.sub(r"[_-]+", " ", value)
    cleaned = re.sub(r"([a-z])([A-Z])", r"\1 \2", cleaned)
    return cleaned.strip().capitalize()


def _map_alignment(value: Optional[str]) -> PP_ALIGN:
    if value == 'center':
        return PP_ALIGN.CENTER
    if value == 'right':
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _apply_font_formatting(font, formatting: dict) -> None:
    font_name = formatting.get('fontFamily')
    if isinstance(font_name, str) and font_name.strip():
        font.name = font_name.strip()

    font_size = formatting.get('fontSize')
    numeric_size = _safe_float(font_size, 0)
    if numeric_size > 0:
        font.size = Pt(_px_to_pt(numeric_size))

    font.bold = bool(formatting.get('bold'))
    font.italic = bool(formatting.get('italic'))
    font.underline = bool(formatting.get('underline'))
    font.strike = bool(formatting.get('strikethrough'))

    color = _parse_hex_color(formatting.get('color'))
    if color is not None:
        font.color.rgb = color


def _format_table_value(value: Any) -> str:
    if value is None:
        return ''
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        if math.isnan(value) or math.isinf(value):
            return ''
        return f"{value}"
    return str(value)


def _extract_table_preview(metadata: Optional[dict[str, Any]]) -> Optional[dict[str, list[Any]]]:
    if not metadata:
        return None

    candidates = [
        metadata.get('tableData'),
        metadata.get('previewTable'),
        metadata.get('table'),
        metadata.get('data'),
        metadata.get('rows'),
    ]

    for candidate in candidates:
        if not candidate:
            continue

        if isinstance(candidate, dict):
            headers = _ensure_list(candidate.get('headers'))
            rows = _ensure_list(candidate.get('rows'))

            if headers and rows:
                return {'headers': headers, 'rows': rows}

            data_rows = _ensure_list(candidate.get('data'))
            if data_rows and isinstance(data_rows[0], dict):
                keys = list(data_rows[0].keys())
                return {'headers': keys, 'rows': data_rows}

        if isinstance(candidate, list) and candidate and isinstance(candidate[0], dict):
            keys = list(candidate[0].keys())
            return {'headers': keys, 'rows': candidate}

    return None


def _build_table_cells(headers: list[Any], rows: list[Any]) -> list[list[dict[str, Any]]]:
    table: list[list[dict[str, Any]]] = []
    header_cells = [
        {'content': _format_table_value(header), 'formatting': {'bold': True, 'align': 'left'}}
        for header in headers
    ]
    table.append(header_cells)

    for row in rows:
        cells: list[dict[str, Any]] = []
        if isinstance(row, dict):
            for header in headers:
                cells.append({'content': _format_table_value(row.get(header))})
        elif isinstance(row, (list, tuple)):
            sequence = list(row)
            for index, header in enumerate(headers):
                value = sequence[index] if index < len(sequence) else ''
                cells.append({'content': _format_table_value(value)})
        else:
            cells.append({'content': _format_table_value(row)})

        table.append(cells)

    return table


def _is_numeric_value(value: Any) -> bool:
    if isinstance(value, bool):
        return False
    try:
        numeric = float(value)
    except (TypeError, ValueError):
        return False
    return math.isfinite(numeric)


def _derive_numeric_keys(sample: dict[str, Any]) -> list[str]:
    return [key for key, value in sample.items() if _is_numeric_value(value)]


def _derive_category_key(sample: dict[str, Any]) -> str:
    preferred = ['label', 'name', 'category', 'x', 'dimension']
    for key in preferred:
        if key in sample:
            return key
    for key, value in sample.items():
        if isinstance(value, str):
            return key
    return 'index'


def _extract_chart_colors(candidate: Any) -> list[str]:
    colors = []
    if isinstance(candidate, dict):
        options = [
            candidate.get('chartColors'),
            candidate.get('chart_colors'),
            candidate.get('colorPalette'),
            candidate.get('color_palette'),
            candidate.get('colors'),
        ]
        for option in options:
            if isinstance(option, list):
                colors.extend([color for color in option if _is_non_empty_str(color)])
    elif isinstance(candidate, list):
        colors.extend([color for color in candidate if _is_non_empty_str(color)])
    return colors


def _normalise_series_colors(metadata: dict[str, Any], candidate: Optional[dict[str, Any]]) -> list[str]:
    colors: list[str] = []
    if candidate:
        colors.extend(_extract_chart_colors(candidate))
    colors.extend(_extract_chart_colors(metadata))
    seen: set[str] = set()
    unique: list[str] = []
    for color in colors:
        key = color.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(color)
    return unique


def _extract_chart_preview(metadata: Optional[dict[str, Any]]) -> Optional[dict[str, Any]]:
    if not metadata:
        return None

    # Check chartContext first (Chart Maker stores data here)
    chart_context = _as_dict(metadata.get('chartContext') or metadata.get('chart_context'))
    chart_config = _as_dict(chart_context.get('chartConfig') or chart_context.get('chart_config')) if chart_context else None
    context_data = chart_config.get('data') if chart_config else None
    
    # Also check for chartState (chart configuration)
    chart_state = _as_dict(metadata.get('chartState') or metadata.get('chart_state'))
    
    # Log what we found for debugging
    if chart_context:
        logger.info('Found chartContext with keys: %s', list(chart_context.keys())[:10])
    if chart_config:
        logger.info('Found chartConfig with keys: %s, has data: %s', 
                   list(chart_config.keys())[:10], bool(context_data))
    if chart_state:
        logger.info('Found chartState with keys: %s, chartType: %s', 
                   list(chart_state.keys())[:10], chart_state.get('chartType') or chart_state.get('chart_type'))
    
    # Try multiple sources for the actual data
    candidate_raw = (
        metadata.get('chartData')
        or metadata.get('chart_data')
        or metadata.get('chart')
        or metadata.get('chartMetadata')
        or metadata.get('chart_metadata')
        or context_data  # Chart Maker data location
        or metadata.get('visualisation')
        or metadata.get('visualization')
    )

    candidate_dict = _as_dict(candidate_raw)
    raw_source: list[Any] = []
    if isinstance(candidate_raw, list):
        raw_source = candidate_raw
    elif candidate_dict and isinstance(candidate_dict.get('data'), list):
        raw_source = candidate_dict['data']
    elif isinstance(metadata.get('data'), list):
        raw_source = metadata['data']

    records = [entry for entry in raw_source if isinstance(entry, dict)]
    if not records:
        logger.debug('No chart records found in metadata. chartContext: %s, chartState: %s', 
                    bool(chart_context), bool(chart_state))
        return None

    sample = records[0]
    
    # Extract chart type from multiple sources (prioritize chartState)
    chart_type_raw = (
        (chart_state.get('chartType') if chart_state else None)
        or (chart_state.get('chart_type') if chart_state else None)
        or (chart_config.get('chartType') if chart_config else None)
        or (chart_config.get('chart_type') if chart_config else None)
        or (candidate_dict.get('chart_type') if candidate_dict else None)
        or (candidate_dict.get('type') if candidate_dict else None)
        or metadata.get('chartType')
        or metadata.get('chart_type')
        or metadata.get('type')
    )
    chart_type = str(chart_type_raw).lower() if chart_type_raw else ''

    numeric_keys = _derive_numeric_keys(sample)

    label_key = _derive_category_key(sample)
    if label_key not in sample:
        label_key = 'index'

    possible_pie = False
    if 'value' in sample and _is_numeric_value(sample['value']):
        possible_pie = 'label' in sample or 'name' in sample or 'x' in sample
    elif 'y' in sample and _is_numeric_value(sample['y']):
        possible_pie = 'label' in sample or 'name' in sample or 'x' in sample

    forced_pie = any(token in chart_type for token in ('pie', 'donut', 'doughnut'))

    if forced_pie or (not chart_type and possible_pie):
        value_key = 'value' if 'value' in sample else 'y' if 'y' in sample else None
        if not value_key and numeric_keys:
            value_key = numeric_keys[0]
        if not value_key:
            return None

        categories: list[str] = []
        values: list[float] = []
        for index, entry in enumerate(records):
            if not isinstance(entry, dict):
                continue
            raw_value = entry.get(value_key)
            numeric = _safe_float(raw_value, math.nan)
            if math.isnan(numeric):
                continue
            label = entry.get(label_key)
            if label is None:
                label = f'Slice {index + 1}'
            categories.append(str(label))
            values.append(numeric)

        if not values:
            return None

        chart_kind = 'donut' if forced_pie and ('donut' in chart_type or 'doughnut' in chart_type) else 'pie'
        return {
            'type': chart_kind,
            'categories': categories,
            'series': [{'name': 'Series 1', 'values': values}],
            'seriesColors': _normalise_series_colors(metadata, candidate_dict),
        }

    if not numeric_keys:
        return None

    chart_kind = 'column'
    is_stacked_bar = 'stacked' in chart_type and 'bar' in chart_type
    if 'line' in chart_type:
        chart_kind = 'line'
    elif 'area' in chart_type:
        chart_kind = 'area'
    elif 'scatter' in chart_type:
        chart_kind = 'scatter'
    elif 'bar' in chart_type:
        chart_kind = 'bar' if not is_stacked_bar else 'stacked_bar'

    # For stacked bar charts, check if data is in long format (needs pivoting)
    # Long format: each row is a category-segment combination
    # Wide format: each row is a category with multiple segment columns
    legend_field = (
        (chart_state.get('legendField') if chart_state else None)
        or (chart_state.get('legend_field') if chart_state else None)
        or metadata.get('legendField')
        or metadata.get('legend_field')
    )
    
    categories: list[str] = []
    series: list[dict[str, Any]] = []
    
    # If we have a legend field and it's a stacked bar, we need to pivot the data
    if is_stacked_bar and legend_field and label_key:
        # Check if data is in long format (legend_field values appear as rows)
        # Look for legend_field values in the records
        legend_values = set()
        for entry in records:
            if isinstance(entry, dict) and legend_field in entry:
                legend_value = entry.get(legend_field)
                if legend_value is not None:
                    legend_values.add(str(legend_value))
        
        if legend_values and len(legend_values) > 1:
            # Data is in long format - need to pivot
            # Group by category (label_key) and create series for each legend value
            grouped_data: dict[str, dict[str, float]] = {}  # {category: {legend_value: numeric_value}}
            all_categories_set = set()
            
            for entry in records:
                if not isinstance(entry, dict):
                    continue
                category = str(entry.get(label_key, ''))
                legend_value = str(entry.get(legend_field, ''))
                
                if not category or not legend_value:
                    continue
                
                all_categories_set.add(category)
                
                if category not in grouped_data:
                    grouped_data[category] = {}
                
                # Get the numeric value from the first numeric key found
                for key in numeric_keys:
                    numeric = _safe_float(entry.get(key), math.nan)
                    if not math.isnan(numeric):
                        # Use legend value as series name, accumulate if multiple numeric keys
                        if legend_value not in grouped_data[category]:
                            grouped_data[category][legend_value] = 0.0
                        grouped_data[category][legend_value] += numeric
                        break
            
            # Create series for each legend value
            categories = sorted(all_categories_set)  # Sort for consistent ordering
            sorted_legend_values = sorted(legend_values)
            
            for legend_val in sorted_legend_values:
                values = []
                for cat in categories:
                    value = grouped_data.get(cat, {}).get(legend_val, 0.0)
                    values.append(value)
                
                # Clean up series name
                clean_name = str(legend_val)
                if '_trace_' in clean_name:
                    clean_name = re.sub(r'_trace_\d+$', '', clean_name)
                clean_name = clean_name.rstrip('_')
                series.append({'name': clean_name, 'values': values})
            
            logger.info('Stacked bar chart: Pivoted data from long format. Categories: %d, Series: %d, Legend field: %s', 
                       len(categories), len(series), legend_field)
        else:
            # Data is already in wide format, use existing logic
            series_values: dict[str, list[float]] = {key: [] for key in numeric_keys}
            for index, entry in enumerate(records):
                if not isinstance(entry, dict):
                    continue
                label = entry.get(label_key, f'Row {index + 1}')
                categories.append(str(label))
                for key in numeric_keys:
                    numeric = _safe_float(entry.get(key), math.nan)
                    if math.isnan(numeric):
                        numeric = 0.0
                    series_values[key].append(numeric)
            
            for key in numeric_keys:
                # Clean up series name: remove "_trace_X" suffix for cleaner legend
                clean_name = str(key)
                if '_trace_' in clean_name:
                    # Remove "_trace_0", "_trace_1", etc.
                    clean_name = re.sub(r'_trace_\d+$', '', clean_name)
                # Also remove any trailing underscores
                clean_name = clean_name.rstrip('_')
                series.append({'name': clean_name, 'values': series_values[key]})
    else:
        # Not a stacked bar or no legend field - use existing logic
        # Special handling for scatter plots: extract X and Y coordinate pairs
        if chart_kind == 'scatter':
            # For scatter plots, we need X and Y pairs
            # Check if data has 'x' and 'y' keys, or if we need to use first two numeric keys
            x_key = None
            y_key = None
            
            # First, try to find explicit 'x' and 'y' keys
            if 'x' in sample and 'y' in sample:
                x_key = 'x'
                y_key = 'y'
            elif 'X' in sample and 'Y' in sample:
                x_key = 'X'
                y_key = 'Y'
            elif numeric_keys and len(numeric_keys) >= 2:
                # Use first two numeric keys as X and Y
                x_key = numeric_keys[0]
                y_key = numeric_keys[1]
            else:
                logger.warning('Scatter plot: Unable to find X and Y keys in data. Sample keys: %s', list(sample.keys()))
                return None
            
            # Extract X and Y values as numeric pairs
            x_values: list[float] = []
            y_values: list[float] = []
            for entry in records:
                if not isinstance(entry, dict):
                    continue
                x_val = _safe_float(entry.get(x_key), math.nan)
                y_val = _safe_float(entry.get(y_key), math.nan)
                if not math.isnan(x_val) and not math.isnan(y_val):
                    x_values.append(x_val)
                    y_values.append(y_val)
            
            if not x_values or not y_values:
                logger.warning('Scatter plot: No valid X/Y pairs found. X values: %d, Y values: %d', len(x_values), len(y_values))
                return None
            
            # For scatter plots, we create a single series with X/Y pairs
            # Store as special format: series with 'x_values' and 'y_values' keys
            series.append({
                'name': 'Scatter Series',
                'x_values': x_values,
                'y_values': y_values
            })
            
            # For scatter plots, categories are the X values (for axis labels)
            categories = [str(x) for x in x_values]
            
            logger.info('Scatter plot: Extracted %d X/Y pairs. X range: [%.2f, %.2f], Y range: [%.2f, %.2f]', 
                       len(x_values), min(x_values), max(x_values), min(y_values), max(y_values))
        else:
            # Regular chart logic (bar, line, etc.)
            series_values: dict[str, list[float]] = {key: [] for key in numeric_keys}
            for index, entry in enumerate(records):
                if not isinstance(entry, dict):
                    continue
                label = entry.get(label_key, f'Row {index + 1}')
                categories.append(str(label))
                for key in numeric_keys:
                    numeric = _safe_float(entry.get(key), math.nan)
                    if math.isnan(numeric):
                        numeric = 0.0
                    series_values[key].append(numeric)
            
            for key in numeric_keys:
                # Clean up series name: remove "_trace_X" suffix for cleaner legend
                clean_name = str(key)
                if '_trace_' in clean_name:
                    # Remove "_trace_0", "_trace_1", etc.
                    clean_name = re.sub(r'_trace_\d+$', '', clean_name)
                # Also remove any trailing underscores
                clean_name = clean_name.rstrip('_')
                series.append({'name': clean_name, 'values': series_values[key]})

    # Extract chart configuration from multiple sources
    legend_position = (
        (chart_state.get('legendPosition') if chart_state else None)
        or (chart_state.get('legend_position') if chart_state else None)
        or (chart_config.get('legendPosition') if chart_config else None)
        or metadata.get('legendPosition')
        or metadata.get('legend_position')
    )
    
    show_legend = (
        (chart_state.get('showLegend') if chart_state else None)
        or (chart_state.get('show_legend') if chart_state else None)
        or (chart_config.get('showLegend') if chart_config else None)
        or (chart_config.get('show_legend') if chart_config else None)
        or metadata.get('showLegend')
        or metadata.get('show_legend')
    )
    
    x_axis_label = (
        (chart_state.get('xAxisLabel') if chart_state else None)
        or (chart_state.get('x_axis_label') if chart_state else None)
        or (chart_state.get('xAxis') if chart_state else None)
        or (chart_config.get('xAxisLabel') if chart_config else None)
        or metadata.get('xAxisLabel')
        or metadata.get('x_axis_label')
    )
    
    y_axis_label = (
        (chart_state.get('yAxisLabel') if chart_state else None)
        or (chart_state.get('y_axis_label') if chart_state else None)
        or (chart_state.get('yAxis') if chart_state else None)
        or (chart_config.get('yAxisLabel') if chart_config else None)
        or metadata.get('yAxisLabel')
        or metadata.get('y_axis_label')
    )
    
    show_values = (
        (chart_state.get('showDataLabels') if chart_state else None)
        or (chart_state.get('show_data_labels') if chart_state else None)
        or (chart_config.get('showDataLabels') if chart_config else None)
        or (chart_config.get('show_data_labels') if chart_config else None)
        or metadata.get('showValues')
        or metadata.get('showDataLabels')
    )
    
    axis_zero = metadata.get('axisIncludesZero') or metadata.get('includeZero')
    
    # Extract color palette
    color_palette = (
        (chart_state.get('colorPalette') if chart_state else None)
        or (chart_state.get('color_palette') if chart_state else None)
        or (chart_config.get('colors') if chart_config else None)
        or (chart_config.get('colorPalette') if chart_config else None)
    )

    # Normalize series colors from multiple sources
    series_colors = _normalise_series_colors(metadata, candidate_dict)
    if color_palette and isinstance(color_palette, list):
        # Prepend color palette from chartState/chartConfig
        series_colors = [str(c) for c in color_palette if _is_non_empty_str(str(c))] + series_colors
    
    # Extract dual axis configuration
    second_y_axis = (
        (chart_state.get('secondYAxis') if chart_state else None)
        or (chart_state.get('second_y_axis') if chart_state else None)
        or metadata.get('secondYAxis')
        or metadata.get('second_y_axis')
    )
    
    dual_axis_mode = (
        (chart_state.get('dualAxisMode') if chart_state else None)
        or (chart_state.get('dual_axis_mode') if chart_state else None)
        or metadata.get('dualAxisMode')
        or metadata.get('dual_axis_mode')
    )
    
    y_axis_field = (
        (chart_state.get('yAxis') if chart_state else None)
        or (chart_state.get('y_axis') if chart_state else None)
        or metadata.get('yAxis')
        or metadata.get('y_axis')
    )
    
    return {
        'type': chart_kind,
        'categories': categories,
        'series': series,
        'legendPosition': legend_position,
        'showLegend': bool(show_legend) if show_legend is not None else None,
        'showValues': bool(show_values) if show_values is not None else False,
        'axisIncludesZero': bool(axis_zero) if axis_zero is not None else False,
        'xAxisLabel': str(x_axis_label) if x_axis_label else None,
        'yAxisLabel': str(y_axis_label) if y_axis_label else None,
        'yAxis': str(y_axis_field) if y_axis_field else None,
        'secondYAxis': str(second_y_axis) if second_y_axis else None,
        'dualAxisMode': str(dual_axis_mode) if dual_axis_mode else None,
        'seriesColors': series_colors,
    }


def _extract_summary_lines(metadata: Optional[dict[str, Any]]) -> list[str]:
    if not metadata:
        return []

    lines: list[str] = []
    summary = metadata.get('summary')
    if isinstance(summary, list):
        lines.extend(str(item) for item in summary if _is_non_empty_str(item))
    elif _is_non_empty_str(summary):
        lines.append(str(summary))

    keys_to_skip = {
        'chartData',
        'chart_data',
        'chart',
        'chartMetadata',
        'chart_metadata',
        'chartState',
        'chart_state',
        'chartContext',
        'chart_context',
        'chartId',
        'chart_id',
        'chartTitle',
        'chart_title',
        'capturedAt',
        'captured_at',
        'sourceAtomTitle',
        'source_atom_title',
        'previewTable',
        'tableData',
        'table',
        'rows',
        'data',
        'previewImage',
        'preview_image',
        'image',
        'thumbnail',
        'visualization',
        'visualisation',
        'visualizationManifest',
        'visualisationManifest',
        'manifest',
    }

    for key, value in metadata.items():
        if key in keys_to_skip:
            continue
        if isinstance(value, (str, int, float, bool)):
            lines.append(f"{_humanise_key(key)}: {_format_table_value(value)}")
        elif isinstance(value, list) and value and len(lines) < 8:
            preview = ', '.join(str(item) for item in value[:3])
            lines.append(f"{_humanise_key(key)}: {preview}")
        if len(lines) >= 8:
            break

    return lines[:8]


def _normalise_card_layout(value: Optional[str]) -> str:
    if not value:
        return 'none'
    lowered = value.lower()
    if lowered in {'none', 'top', 'bottom', 'left', 'right', 'full'}:
        return lowered
    return 'none'


def _resolve_overlay_fill(settings: dict) -> tuple[str, object]:
    accent_image = settings.get('accentImage') or settings.get('accent_image')
    if isinstance(accent_image, str) and accent_image.strip():
        return 'image', accent_image.strip()

    card_color = settings.get('cardColor') or settings.get('card_color')
    if isinstance(card_color, str) and card_color.strip():
        token = card_color.strip()
        if token.startswith('solid-') and len(token) >= 12:
            return 'solid', f"#{token[6:12]}"

        lookup = token.lower()
        preset = GRADIENT_PRESETS.get(lookup)
        if preset:
            return 'gradient', preset

    return 'solid', DEFAULT_OVERLAY_COLOR


def _compute_overlay_rect(layout: str, width: float, height: float) -> Optional[tuple[float, float, float, float]]:
    if width <= 0 or height <= 0:
        return None

    if layout == 'full':
        return 0.0, 0.0, width, height

    if layout == 'top':
        ratio = TOP_LAYOUT_MIN_HEIGHT / CANVAS_STAGE_HEIGHT if CANVAS_STAGE_HEIGHT else 0
        overlay_height = max(TOP_LAYOUT_MIN_HEIGHT, height * ratio)
        overlay_height = min(overlay_height, height)
        return 0.0, 0.0, width, overlay_height

    if layout == 'bottom':
        ratio = BOTTOM_LAYOUT_MIN_HEIGHT / CANVAS_STAGE_HEIGHT if CANVAS_STAGE_HEIGHT else 0
        overlay_height = max(BOTTOM_LAYOUT_MIN_HEIGHT, height * ratio)
        overlay_height = min(overlay_height, height)
        return 0.0, height - overlay_height, width, overlay_height

    if layout == 'left':
        overlay_width = max(SIDE_LAYOUT_MIN_WIDTH, width * SIDE_LAYOUT_RATIO)
        overlay_width = min(overlay_width, width)
        return 0.0, 0.0, overlay_width, height

    if layout == 'right':
        overlay_width = max(SIDE_LAYOUT_MIN_WIDTH, width * SIDE_LAYOUT_RATIO)
        overlay_width = min(overlay_width, width)
        return width - overlay_width, 0.0, overlay_width, height

    return None


def _render_layout_overlay(
    slide,
    slide_payload: SlideExportPayload,
    base_width: float,
    base_height: float,
    offset_x: float = 0.0,
    offset_y: float = 0.0,
) -> None:
    settings = slide_payload.presentation_settings or {}
    if not isinstance(settings, dict):
        return

    layout_value = settings.get('cardLayout') or settings.get('card_layout')
    layout = _normalise_card_layout(layout_value)
    if layout == 'none':
        return

    rect = _compute_overlay_rect(layout, base_width, base_height)
    if not rect:
        return

    x, y, width, height = rect
    fill_type, fill_value = _resolve_overlay_fill(settings)

    if fill_type == 'image' and isinstance(fill_value, str):
        try:
            image_bytes = _decode_data_url(fill_value)
            image_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(
                image_stream,
                _px_to_emu(x + offset_x),
                _px_to_emu(y + offset_y),
                width=_px_to_emu(width),
                height=_px_to_emu(height),
            )
            return
        except ExportGenerationError:
            fill_type = 'solid'
            fill_value = DEFAULT_OVERLAY_COLOR
        except Exception as exc:  # pragma: no cover - best effort logging
            logger.warning('Unable to render accent image for slide %s: %s', slide_payload.id, exc)
            fill_type = 'solid'
            fill_value = DEFAULT_OVERLAY_COLOR

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _px_to_emu(x + offset_x),
        _px_to_emu(y + offset_y),
        _px_to_emu(width),
        _px_to_emu(height),
    )
    shape.line.fill.background()

    if fill_type == 'gradient' and isinstance(fill_value, dict):
        stops = fill_value.get('stops') if isinstance(fill_value.get('stops'), list) else []
        colors = [color for color in stops if isinstance(color, str)]
        if len(colors) >= 2:
            gradient = shape.fill
            gradient.gradient()
            try:
                gradient.gradient_angle = int(fill_value.get('angle')) if fill_value.get('angle') is not None else 135
            except (TypeError, ValueError):
                gradient.gradient_angle = 135

            start_color = _parse_hex_color(colors[0]) or _parse_hex_color(DEFAULT_OVERLAY_COLOR)
            end_color = _parse_hex_color(colors[-1]) or start_color
            if start_color is not None and end_color is not None:
                gradient_stops = gradient.gradient_stops
                gradient_stops[0].position = 0.0
                gradient_stops[0].color.rgb = start_color
                gradient_stops[1].position = 1.0
                gradient_stops[1].color.rgb = end_color
                return

        fill_type = 'solid'
        fill_value = colors[0] if colors else DEFAULT_OVERLAY_COLOR

    if fill_type == 'solid':
        rgb = _parse_hex_color(str(fill_value)) or _parse_hex_color(DEFAULT_OVERLAY_COLOR)
        if rgb is not None:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = rgb


def _render_text_box(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0) -> None:
    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        return

    shape = slide.shapes.add_textbox(
        _px_to_emu(obj.x + offset_x),
        _px_to_emu(obj.y + offset_y),
        _px_to_emu(width),
        _px_to_emu(height),
    )

    text = _html_to_plain_text(str(obj.props.get('text', '') or ''))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    formatting = obj.props
    lines = text.split('\n') if text else ['']
    for index, line in enumerate(lines):
        paragraph = text_frame.add_paragraph() if index > 0 else text_frame.paragraphs[0]
        paragraph.text = line
        paragraph.alignment = _map_alignment(formatting.get('align'))
        _apply_font_formatting(paragraph.font, formatting)

    rotation = _safe_float(obj.rotation, 0)
    if rotation:
        shape.rotation = rotation


def _render_image(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0) -> None:
    """
    Render an image object on a PowerPoint slide.
    
    Handles images from:
    - Data URLs (base64 encoded)
    - Remote URLs (HTTP/HTTPS)
    - Various image formats (converts to PNG/JPEG as needed)
    
    Args:
        slide: PowerPoint slide object
        obj: Image object payload with props containing 'src' or 'dataUrl'
        offset_x: X offset in pixels
        offset_y: Y offset in pixels
    """
    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        logger.warning('Skipping image %s: invalid dimensions (width=%s, height=%s)', obj.id, width, height)
        return

    source = obj.props.get('src') or obj.props.get('dataUrl') or obj.props.get('data_url')
    if not isinstance(source, str) or not source:
        logger.warning('Skipping image %s: missing source (src/dataUrl). Props keys: %s', 
                     obj.id, list(obj.props.keys()) if obj.props else 'none')
        return
    
    logger.debug('PPTX: Processing image %s with source type: %s (length: %d)', 
                obj.id, 'URL' if source.startswith(('http://', 'https://')) else 'data URL' if source.startswith('data:') else 'unknown',
                len(source) if isinstance(source, str) else 0)

    try:
        # Use robust image loading helper
        image_bytes, format_hint = _load_image_asset(source, obj_id=obj.id)
        
        if not image_bytes:
            logger.error('Image %s: loaded but got empty bytes', obj.id)
            return
        
        image_stream = io.BytesIO(image_bytes)
        
        logger.debug(
            'Adding image %s to slide: %dx%d at (%.1f, %.1f), format: %s, size: %d bytes',
            obj.id, width, height, obj.x + offset_x, obj.y + offset_y, format_hint or 'unknown', len(image_bytes)
        )
        
        try:
            shape = slide.shapes.add_picture(
                image_stream,
                _px_to_emu(obj.x + offset_x),
                _px_to_emu(obj.y + offset_y),
                width=_px_to_emu(width),
                height=_px_to_emu(height),
            )
            
            # Set rotation if specified
            rotation = _safe_float(obj.rotation, 0)
            if rotation:
                shape.rotation = rotation
            
            logger.info('Successfully rendered image %s on slide', obj.id)
        except Exception as add_picture_exc:
            # python-pptx add_picture can fail for various reasons (invalid format, corrupted data, etc.)
            logger.warning('Failed to add image %s to PowerPoint slide using add_picture: %s. Skipping.', 
                         obj.id, add_picture_exc)
            return  # Skip this image and continue
        
    except ExportGenerationError as exc:
        # Log error but skip this image instead of failing entire export
        logger.warning('Failed to load image %s: %s. Skipping this image.', obj.id, exc)
        return  # Gracefully skip this image and continue with export
    except Exception as exc:
        # Log unexpected errors but don't fail entire export
        logger.warning('Unexpected error rendering image %s: %s. Skipping this image.', obj.id, exc)
        return  # Gracefully skip this image and continue with export


def _map_dash_style(style: Optional[str]) -> Optional[MSO_LINE_DASH_STYLE]:
    if not style:
        return None

    lookup = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "dotted": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dash-dot": MSO_LINE_DASH_STYLE.DASH_DOT,
    }

    return lookup.get(style.lower())


def _apply_shape_styles(shape, props: dict, *, is_line: bool = False) -> None:
    opacity = _safe_float(props.get("opacity"), 1.0)
    opacity = min(max(opacity, 0.0), 1.0)

    if not is_line:
        fill_color = _parse_color_token(props.get("fill"))
        if fill_color is None or opacity <= 0:
            shape.fill.background()
        else:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = fill_color
            if opacity < 1.0:
                fill.fore_color.transparency = max(0.0, min(1.0, 1.0 - opacity))

    stroke_width = _safe_float(props.get("strokeWidth"), 0.0)
    stroke_color = _parse_color_token(props.get("stroke"))
    line = getattr(shape, "line", None)

    if line is not None:
        if stroke_color is None or stroke_width <= 0:
            try:
                line.fill.background()
            except AttributeError:
                pass
            line.width = 0
        else:
            try:
                line.color.rgb = stroke_color
            except AttributeError:
                pass
            line.width = Pt(_px_to_pt(stroke_width))
            dash = _map_dash_style(props.get("strokeStyle"))
            if dash is not None:
                line.dash_style = dash
            if opacity < 1.0:
                try:
                    line.fill.solid()
                    line.fill.fore_color.rgb = stroke_color
                    line.fill.fore_color.transparency = max(0.0, min(1.0, 1.0 - opacity))
                except AttributeError:
                    pass


def _resolve_line_points(
    shape_id: str,
    obj: SlideExportObjectPayload,
    offset_x: float,
    offset_y: float,
) -> tuple[float, float, float, float]:
    width = _safe_float(obj.width, 0.0)
    height = _safe_float(obj.height, 0.0)
    left = obj.x + offset_x
    top = obj.y + offset_y

    if shape_id == "line-vertical":
        x = left + (width / 2.0)
        return x, top, x, top + height

    if shape_id == "line-diagonal":
        return left, top + height, left + width, top

    # Default to horizontal line
    y = top + (height / 2.0)
    return left, y, left + width, y


def _render_shape(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0) -> None:
    props = obj.props or {}
    shape_id = str(props.get("shapeId") or props.get("shape_id") or "").strip()

    if shape_id in LINE_SHAPES:
        x1, y1, x2, y2 = _resolve_line_points(shape_id, obj, offset_x, offset_y)
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _px_to_emu(x1),
            _px_to_emu(y1),
            _px_to_emu(x2),
            _px_to_emu(y2),
        )
        _apply_shape_styles(shape, props, is_line=True)
    else:
        width = _safe_float(obj.width, 0.0)
        height = _safe_float(obj.height, 0.0)
        if width <= 0 or height <= 0:
            return

        mapped_shape = SHAPE_ID_TO_MSO.get(shape_id, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(
            mapped_shape,
            _px_to_emu(obj.x + offset_x),
            _px_to_emu(obj.y + offset_y),
            _px_to_emu(width),
            _px_to_emu(height),
        )
        _apply_shape_styles(shape, props)

    rotation = _safe_float(obj.rotation, 0.0)
    if rotation:
        shape.rotation = rotation

def _extract_table_data(obj: SlideExportObjectPayload) -> Optional[list]:
    data = obj.props.get('data')
    if isinstance(data, list) and data:
        return data
    return None


def _apply_table_cell(cell, payload: dict) -> None:
    text = _html_to_plain_text(str(payload.get('content', '') or ''))
    formatting = payload.get('formatting') or {}

    cell.text = ''
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = cell.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = _map_alignment(formatting.get('align'))
    _apply_font_formatting(paragraph.font, formatting)


def _render_table(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0) -> None:
    table_data = _extract_table_data(obj)
    if not table_data:
        return

    rows = len(table_data)
    cols = len(table_data[0]) if rows else 0
    if rows == 0 or cols == 0:
        return

    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        return

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        _px_to_emu(obj.x + offset_x),
        _px_to_emu(obj.y + offset_y),
        _px_to_emu(width),
        _px_to_emu(height),
    )

    table = table_shape.table

    column_width = _px_to_emu(width) // cols
    for index in range(cols):
        table.columns[index].width = column_width

    row_height = _px_to_emu(height) // rows
    for index in range(rows):
        table.rows[index].height = row_height

    for row_index, row in enumerate(table_data):
        for col_index in range(cols):
            try:
                payload = row[col_index]
            except IndexError:
                payload = ''

            cell = table.cell(row_index, col_index)
            if isinstance(payload, dict):
                _apply_table_cell(cell, payload)
            else:
                cell.text = ''
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.text = str(payload or '')
                paragraph.alignment = PP_ALIGN.LEFT

    rotation = _safe_float(obj.rotation, 0)
    if rotation:
        table_shape.rotation = rotation


def _map_chart_type(chart_type: Optional[str]) -> XL_CHART_TYPE:
    chart_type_lower = (chart_type or '').lower()
    mapping = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,  # Bar charts render as vertical columns (not horizontal bars)
        'stacked_bar': XL_CHART_TYPE.COLUMN_STACKED,
        'stacked-bar': XL_CHART_TYPE.COLUMN_STACKED,
        'stackedbar': XL_CHART_TYPE.COLUMN_STACKED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'donut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'scatter_lines': XL_CHART_TYPE.XY_SCATTER_LINES,
    }
    return mapping.get(chart_type_lower, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _map_legend_position(position: Optional[str]) -> Optional[XL_LEGEND_POSITION]:
    if not position:
        return None
    lookup = {
        'top': XL_LEGEND_POSITION.TOP,
        'bottom': XL_LEGEND_POSITION.BOTTOM,
        'left': XL_LEGEND_POSITION.LEFT,
        'right': XL_LEGEND_POSITION.RIGHT,
    }
    return lookup.get(position.lower())


def _resolve_post_animation_image(props: dict[str, Any]) -> Optional[bytes]:
    png_candidate = props.get('postAnimationPng') or props.get('postAnimationImage')
    if isinstance(png_candidate, str) and png_candidate.strip():
        try:
            return _decode_data_url(png_candidate)
        except ExportGenerationError as exc:
            logger.warning('Unable to decode post-animation PNG for chart: %s', exc)

    svg_candidate = props.get('postAnimationSvg')
    if isinstance(svg_candidate, str) and svg_candidate.strip():
        if cairosvg is None:
            logger.debug('SVG chart provided but cairosvg is unavailable; skipping rasterisation.')
        else:
            try:
                svg_bytes = _decode_svg_data_url(svg_candidate)
                return cairosvg.svg2png(bytestring=svg_bytes)
            except ExportGenerationError as exc:
                logger.warning('Unable to decode post-animation SVG for chart: %s', exc)
            except Exception as exc:  # pragma: no cover - cairosvg runtime issues
                logger.warning('Failed to rasterise SVG chart for export: %s', exc)

    return None


def _apply_post_animation_overlays(slide: SlideExportPayload) -> None:
    screenshot = getattr(slide, 'screenshot', None)
    if not isinstance(screenshot, SlideScreenshotPayload):
        return

    overlays: list[tuple[SlideExportObjectPayload, bytes]] = []
    for obj in slide.objects:
        props = obj.props or {}
        if not isinstance(props, dict):
            continue
        image_bytes = _resolve_post_animation_image(props)
        if image_bytes:
            overlays.append((obj, image_bytes))

    if not overlays:
        return

    try:
        screenshot_bytes = _decode_data_url(screenshot.data_url)
    except ExportGenerationError as exc:
        logger.warning('Unable to decode screenshot for slide %s: %s', slide.id, exc)
        return

    try:
        base_image = Image.open(io.BytesIO(screenshot_bytes)).convert('RGBA')
    except Exception as exc:  # pragma: no cover - Pillow decoding edge cases
        logger.warning('Unable to load screenshot image for slide %s: %s', slide.id, exc)
        return

    css_width = _safe_float(getattr(screenshot, 'css_width', None), 0)
    css_height = _safe_float(getattr(screenshot, 'css_height', None), 0)
    pixel_ratio = _safe_float(getattr(screenshot, 'pixel_ratio', None), 0)
    image_width = base_image.width
    image_height = base_image.height

    if css_width <= 0 and image_width > 0 and pixel_ratio > 0:
        css_width = image_width / pixel_ratio
    if css_height <= 0 and image_height > 0 and pixel_ratio > 0:
        css_height = image_height / pixel_ratio

    if css_width <= 0:
        css_width = _safe_float(slide.base_width, image_width)
    if css_height <= 0:
        css_height = _safe_float(slide.base_height, image_height)

    scale_x = image_width / css_width if css_width > 0 else 1.0
    scale_y = image_height / css_height if css_height > 0 else 1.0

    updated = False
    for obj, image_bytes in overlays:
        obj_width = _safe_float(obj.width, 0)
        obj_height = _safe_float(obj.height, 0)
        if obj_width <= 0 or obj_height <= 0:
            continue

        try:
            overlay_image = Image.open(io.BytesIO(image_bytes)).convert('RGBA')
        except Exception as exc:  # pragma: no cover - Pillow decoding edge cases
            logger.warning('Unable to decode post-animation image for object %s: %s', obj.id, exc)
            continue

        target_width = max(int(round(obj_width * scale_x)), 1)
        target_height = max(int(round(obj_height * scale_y)), 1)

        if overlay_image.size != (target_width, target_height):
            try:
                overlay_image = overlay_image.resize((target_width, target_height), Image.LANCZOS)
            except Exception:  # pragma: no cover - resize issues
                overlay_image = overlay_image.resize((target_width, target_height))

        dest_x = int(round(_safe_float(obj.x, 0) * scale_x))
        dest_y = int(round(_safe_float(obj.y, 0) * scale_y))

        try:
            base_image.paste(overlay_image, (dest_x, dest_y), overlay_image)
            updated = True
        except Exception as exc:  # pragma: no cover - paste errors
            logger.warning(
                'Unable to overlay post-animation image for object %s on slide %s: %s',
                obj.id,
                slide.id,
                exc,
            )

    if not updated:
        return

    output = io.BytesIO()
    try:
        base_image.save(output, format='PNG')
    except Exception as exc:  # pragma: no cover - save errors
        logger.warning('Unable to encode overlaid screenshot for slide %s: %s', slide.id, exc)
        return

    encoded = base64.b64encode(output.getvalue()).decode('ascii')
    updated_payload = {
        **screenshot.model_dump(by_alias=True),
        'dataUrl': f'data:image/png;base64,{encoded}',
        'width': image_width,
        'height': image_height,
    }

    try:
        slide.screenshot = SlideScreenshotPayload.model_validate(updated_payload)
    except Exception as exc:  # pragma: no cover - validation errors
        logger.warning('Unable to update screenshot payload for slide %s: %s', slide.id, exc)


def _render_chart_image(slide, obj: SlideExportObjectPayload, image_bytes: bytes, offset_x: float, offset_y: float) -> None:
    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        return

    picture = slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        _px_to_emu(obj.x + offset_x),
        _px_to_emu(obj.y + offset_y),
        width=_px_to_emu(width),
        height=_px_to_emu(height),
    )

    rotation = _safe_float(obj.rotation, 0.0)
    if rotation:
        picture.rotation = rotation


def _render_chart(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0) -> None:
    props = obj.props or {}
    
    # Try to render from post-animation image first (if available)
    image_bytes = _resolve_post_animation_image(props)
    if image_bytes:
        logger.debug('Rendering chart %s from post-animation image', obj.id)
        _render_chart_image(slide, obj, image_bytes, offset_x, offset_y)
        return

    # Fall back to rendering from chart data
    data = props.get('chartData') or props.get('chart_data')
    config = props.get('chartConfig') or props.get('chart_config') or {}

    if isinstance(data, list):
        if not data:
            logger.warning('Skipping chart %s: chartData is empty list', obj.id)
            return
    elif not isinstance(data, dict):
        logger.warning('Skipping chart %s: chartData is missing or invalid (got %s)', obj.id, type(data).__name__)
        return
    
    logger.debug('Rendering chart %s from chartData (type: %s)', obj.id, data.get('type') if isinstance(data, dict) else 'unknown')

    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        return

    # Determine chart type early to decide between ChartData and XyChartData
    chart_type_value = data.get('type') if isinstance(data, dict) else config.get('type')
    chart_type = _map_chart_type(chart_type_value)
    is_scatter_chart = chart_type in {XL_CHART_TYPE.XY_SCATTER, XL_CHART_TYPE.XY_SCATTER_LINES, XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS}
    
    series_colors: list[str] = []
    legend_position = _map_legend_position(config.get('legendPosition'))
    
    # Check if legend will be at bottom - reduce chart height to prevent overlap with axis titles
    # The blue background (atom container) contains the chart, so we need to account for legend space
    legend_will_be_bottom = (
        legend_position == XL_LEGEND_POSITION.BOTTOM or
        (legend_position is None and 
         config.get('legendPosition') not in ('right', 'top', 'left'))
    )
    
    # Store original height for positioning calculations
    original_height = height
    
    # Reduce height if legend is at bottom to position legend at the bottom border
    # We want the legend to stick to the bottom edge of the chart border
    if legend_will_be_bottom:
        # Get series count early to adjust reduction based on number of legend items
        # More series = more legend items = need more space
        estimated_series_count = 0
        if isinstance(data, dict):
            series_payload = data.get('series') if isinstance(data.get('series'), list) else []
            estimated_series_count = len(series_payload)
        
        # Adjust reduction based on number of series:
        # - 1-2 series: 30% reduction
        # - 3-4 series: 35% reduction  
        # - 5+ series: 40% reduction
        if estimated_series_count <= 2:
            height_reduction = 0.70  # 30% reduction
        elif estimated_series_count <= 4:
            height_reduction = 0.65  # 35% reduction
        else:
            height_reduction = 0.60  # 40% reduction
        
        height = height * height_reduction
        logger.info('Chart %s: Reduced height by %.0f%% for legend (series count: %d)', 
                   obj.id, (1 - height_reduction) * 100, estimated_series_count)
    
    # Track series count and which series belong to which axis for dual axis support
    series_count = 0
    series_to_axis: dict[int, str] = {}  # Maps series index to 'primary' or 'secondary'
    
    # Check for dual axis configuration
    second_y_axis = config.get('secondYAxis') or config.get('second_y_axis')
    dual_axis_mode = config.get('dualAxisMode') or config.get('dual_axis_mode')
    use_dual_axis = bool(second_y_axis) and dual_axis_mode != 'single'
    
    # Get Y-axis field names to determine which series go on which axis
    y_axis_field = config.get('yAxis') or config.get('y_axis') or data.get('yAxis') if isinstance(data, dict) else None
    second_y_axis_field = config.get('secondYAxis') or config.get('second_y_axis') or data.get('secondYAxis') if isinstance(data, dict) else None
    
    # Log dual axis configuration for debugging
    if use_dual_axis:
        logger.info('Chart %s: Dual axis enabled - yAxis: %s, secondYAxis: %s, mode: %s', 
                   obj.id, y_axis_field, second_y_axis_field, dual_axis_mode)

    # Handle scatter plots with XyChartData
    if is_scatter_chart:
        xy_chart_data = XyChartData()
        
        if isinstance(data, dict):
            series_payload = data.get('series') if isinstance(data.get('series'), list) else []
            
            if not series_payload:
                logger.warning('Scatter plot %s: No series data found', obj.id)
                return
            
            for index, series_entry in enumerate(series_payload):
                if not isinstance(series_entry, dict):
                    continue
                
                name = str(series_entry.get('name') or f'Series {index + 1}')
                # Clean up series name: remove "_trace_X" suffix for cleaner legend
                if '_trace_' in name:
                    name = re.sub(r'_trace_\d+$', '', name)
                name = name.rstrip('_')
                
                # Extract X and Y values for scatter plot
                x_values = series_entry.get('x_values')
                y_values = series_entry.get('y_values')
                
                # If x_values/y_values not found, try to extract from 'values' and 'categories'
                if not x_values or not y_values:
                    # Fallback: use categories as X and values as Y
                    categories = data.get('categories', [])
                    values = series_entry.get('values', [])
                    
                    if categories and values:
                        # Try to convert categories to numeric X values
                        x_values = []
                        for cat in categories:
                            x_val = _safe_float(cat, math.nan)
                            if math.isnan(x_val):
                                # If not numeric, use index
                                x_val = float(len(x_values))
                            x_values.append(x_val)
                        y_values = [_safe_float(v, 0.0) for v in values]
                
                # Validate X and Y values
                if not isinstance(x_values, list) or not isinstance(y_values, list):
                    logger.warning('Scatter plot %s: Series "%s" missing x_values or y_values', obj.id, name)
                    continue
                
                if len(x_values) != len(y_values):
                    logger.warning('Scatter plot %s: Series "%s" has mismatched X/Y lengths (%d vs %d)', 
                                 obj.id, name, len(x_values), len(y_values))
                    # Truncate to minimum length
                    min_len = min(len(x_values), len(y_values))
                    x_values = x_values[:min_len]
                    y_values = y_values[:min_len]
                
                if not x_values or not y_values:
                    logger.warning('Scatter plot %s: Series "%s" has no valid data points', obj.id, name)
                    continue
                
                # Validate all values are numeric
                valid_pairs = []
                for i, (x_val, y_val) in enumerate(zip(x_values, y_values)):
                    x_num = _safe_float(x_val, math.nan)
                    y_num = _safe_float(y_val, math.nan)
                    if not math.isnan(x_num) and not math.isnan(y_num):
                        valid_pairs.append((x_num, y_num))
                    else:
                        logger.debug('Scatter plot %s: Skipping non-numeric point at index %d: x=%s, y=%s', 
                                   obj.id, i, x_val, y_val)
                
                if not valid_pairs:
                    logger.warning('Scatter plot %s: Series "%s" has no valid numeric pairs', obj.id, name)
                    continue
                
                # Log sample data for debugging
                sample_size = min(3, len(valid_pairs))
                logger.info('Scatter plot %s: Series "%s" - %d points. Sample: %s', 
                           obj.id, name, len(valid_pairs), valid_pairs[:sample_size])
                
                # Add series with X/Y pairs
                series = xy_chart_data.add_series(name)
                for x_val, y_val in valid_pairs:
                    series.add_data_point(x_val, y_val)
                
                series_to_axis[series_count] = 'primary'  # Scatter plots typically use primary axis
                series_count += 1
            
            series_colors = [
                str(color)
                for color in _ensure_list(data.get('seriesColors'))
                if _is_non_empty_str(color)
            ]
        else:
            logger.warning('Scatter plot %s: Data is not a dict, cannot extract X/Y pairs', obj.id)
            return
        
        if series_count == 0:
            logger.warning('Scatter plot %s: No valid series created', obj.id)
            return
        
        chart_data = xy_chart_data
        categories = []  # Scatter plots don't use categories
    else:
        # Regular charts (bar, line, etc.) use ChartData
        chart_data = ChartData()
        
        if isinstance(data, dict):
            categories = [str(category) for category in data.get('categories', [])]
            chart_data.categories = categories

            series_payload = data.get('series') if isinstance(data.get('series'), list) else []
            for index, series_entry in enumerate(series_payload):
                if not isinstance(series_entry, dict):
                    continue
                name = str(series_entry.get('name') or f'Series {index + 1}')
                # Clean up series name: remove "_trace_X" suffix for cleaner legend
                if '_trace_' in name:
                    # Remove "_trace_0", "_trace_1", etc.
                    name = re.sub(r'_trace_\d+$', '', name)
                # Also remove any trailing underscores
                name = name.rstrip('_')
                
                # Determine which axis this series should use (for dual axis)
                axis_group = 'primary'
                if use_dual_axis and second_y_axis_field:
                    # Check if this series name matches the second Y-axis field
                    # Series names might be cleaned, so check original name too
                    original_name = str(series_entry.get('name') or '')
                    second_y_axis_lower = str(second_y_axis_field).lower()
                    name_lower = name.lower()
                    original_name_lower = original_name.lower()
                    
                    # Multiple matching strategies:
                    # 1. Exact match or contains match with field name
                    # 2. If we have exactly 2 series, second one goes to secondary axis
                    # 3. Check if series name ends with _trace_1 (second trace)
                    is_secondary = (
                        second_y_axis_lower in name_lower or 
                        second_y_axis_lower in original_name_lower or
                        (len(series_payload) == 2 and index == 1) or  # Second of two series
                        (len(series_payload) > 1 and '_trace_1' in original_name_lower)  # Second trace
                    )
                    
                    if is_secondary:
                        axis_group = 'secondary'
                        logger.info('Chart %s: Series "%s" (index %d, original: "%s") assigned to secondary axis (field: %s)', 
                                   obj.id, name, index, original_name, second_y_axis_field)
                
                values_raw = series_entry.get('values')
                values: list[float] = []
                if isinstance(values_raw, list):
                    for value in values_raw:
                        values.append(_safe_float(value, 0.0))
                if not values and categories:
                    values = [0.0 for _ in categories]
                if values:
                    while len(values) < len(categories):
                        values.append(0.0)
                    if len(values) > len(categories):
                        values = values[: len(categories)]
                    chart_data.add_series(name, values)
                    series_to_axis[series_count] = axis_group
                    series_count += 1
            series_colors = [
                str(color)
                for color in _ensure_list(data.get('seriesColors'))
                if _is_non_empty_str(color)
            ]
        else:
            # Backwards compatibility with legacy chart payloads
            categories = []
            values = []
            for entry in data:
                if isinstance(entry, dict):
                    categories.append(str(entry.get('label', '')))
                    values.append(_safe_float(entry.get('value'), 0))
                else:
                    categories.append(str(entry))
                    values.append(0)
            chart_data.categories = categories
            chart_data.add_series('Series 1', values)
            series_count = 1

    chart_shape = slide.shapes.add_chart(
        chart_type,
        _px_to_emu(obj.x + offset_x),
        _px_to_emu(obj.y + offset_y),
        _px_to_emu(width),
        _px_to_emu(height),
        chart_data,
    )

    chart = chart_shape.chart
    chart.has_title = False

    if legend_position is None and isinstance(data, dict):
        legend_position = _map_legend_position(data.get('legendPosition'))

    # Get axis labels first to determine best legend position
    x_axis_label = config.get('xAxisLabel') or (data.get('xAxisLabel') if isinstance(data, dict) else None)
    y_axis_label = config.get('yAxisLabel') or (data.get('yAxisLabel') if isinstance(data, dict) else None)
    has_axis_labels = bool(x_axis_label or y_axis_label)
    
    # Log axis labels for debugging
    logger.info('Chart %s: Axis labels - X: %s, Y: %s', obj.id, x_axis_label, y_axis_label)
    
    # Always show legend if there are multiple series, or if explicitly requested
    show_legend = config.get('showLegend')
    if show_legend is None and isinstance(data, dict):
        show_legend = data.get('showLegend') or data.get('show_legend')
    
    # Always show legend if there are multiple series, or if explicitly enabled
    # Use series_count from data (more reliable than chart.series which might not be populated yet)
    has_multiple_series = series_count > 1
    
    # For pie charts, always show legend (even with single series) as it's essential
    is_pie_chart = chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}
    
    # Ensure legend is always visible when there are multiple series
    # This fixes the issue where legends don't show on sheet 2
    if show_legend is False:
        # Only hide if explicitly set to False
        chart.has_legend = False
    elif has_multiple_series or is_pie_chart:
        # Always show for multiple series or pie charts (pie charts need legend)
        chart.has_legend = True
        logger.debug('Chart %s: Showing legend (multiple series: %d, is_pie: %s)', obj.id, series_count, is_pie_chart)
    elif show_legend is True:
        # Show if explicitly enabled
        chart.has_legend = True
        logger.debug('Chart %s: Showing legend (explicitly enabled)', obj.id)
    else:
        # Default: show if multiple series or pie chart
        chart.has_legend = has_multiple_series or is_pie_chart
        if chart.has_legend:
            logger.debug('Chart %s: Showing legend (default for %d series, is_pie: %s)', obj.id, series_count, is_pie_chart)
        else:
            logger.debug('Chart %s: Hiding legend (single series: %d)', obj.id, series_count)
    
    if chart.has_legend:
        if legend_position is not None:
            chart.legend.position = legend_position
        else:
            # Use BOTTOM position to place legend at the bottom border of the graph
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        # Keep include_in_layout = True but reduce chart height to position legend at bottom border
        # This ensures the legend sticks to the bottom edge of the chart
        chart.legend.include_in_layout = True
        # Set font size for legend to make it more compact
        try:
            chart.legend.font.size = Pt(9)
        except AttributeError:
            pass
        # Position legend at the very bottom edge of the chart border
        try:
            # With include_in_layout = False, we can position legend more precisely
            # The legend will sit at the bottom border without automatic spacing
            pass
        except AttributeError:
            pass

    plot = chart.plots[0]
    show_values = config.get('showValues')
    if show_values is None and isinstance(data, dict):
        show_values = data.get('showValues')
    
    # Scatter plots don't support has_data_labels the same way as other charts
    # Only set data labels for non-scatter charts
    if not is_scatter_chart:
        if bool(show_values):
            try:
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.number_format = '0.00'
                data_labels.show_value = True
            except (AttributeError, TypeError):
                logger.debug('Chart %s: Could not set data labels (may not be supported for this chart type)', obj.id)
        else:
            try:
                plot.has_data_labels = False
            except (AttributeError, TypeError):
                logger.debug('Chart %s: Could not disable data labels (may not be supported for this chart type)', obj.id)

    axis_includes_zero = config.get('axisIncludesZero')
    if axis_includes_zero is None and isinstance(data, dict):
        axis_includes_zero = data.get('axisIncludesZero')

    # Pie and donut charts don't have axes, handle them separately
    is_pie_or_donut = chart_type in {XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT}
    
    # Scatter plots have different axis structure (both X and Y are value axes)
    is_scatter_plot = chart_type in {XL_CHART_TYPE.XY_SCATTER, XL_CHART_TYPE.XY_SCATTER_LINES, XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS}
    
    if not is_pie_or_donut:
        # Enable and format axes
        try:
            if is_scatter_plot:
                # For scatter plots, both X and Y are value axes
                # X-axis is the category_axis (horizontal axis at bottom)
                # Y-axis is the value_axis (vertical axis on left)
                
                # X-axis (horizontal, bottom) - for scatter plots this is category_axis
                try:
                    category_axis = chart.category_axis
                    category_axis.has_major_gridlines = False  # No vertical gridlines
                    category_axis.has_minor_gridlines = False
                    category_axis.tick_labels.font.size = Pt(9)
                    
                    # Format X-axis numbers as whole numbers (no scientific notation)
                    # Try multiple format approaches for scatter plots
                    try:
                        # First try the standard format
                        category_axis.tick_labels.number_format = '#,##0'
                        logger.debug('Chart %s: Set X-axis number format to #,##0', obj.id)
                    except AttributeError:
                        try:
                            # Alternative format without thousand separator
                            category_axis.tick_labels.number_format = '0'
                            logger.debug('Chart %s: Set X-axis number format to 0', obj.id)
                        except AttributeError:
                            try:
                                # Try General format as fallback
                                category_axis.tick_labels.number_format = 'General'
                                logger.debug('Chart %s: Set X-axis number format to General', obj.id)
                            except AttributeError:
                                logger.warning('Chart %s: Could not set X-axis number format', obj.id)
                    
                    if x_axis_label:
                        category_axis.has_title = True
                        # Ensure axis title exists and set text
                        try:
                            axis_title = category_axis.axis_title
                            text_frame = axis_title.text_frame
                            text_frame.clear()
                            # Ensure at least one paragraph exists
                            if len(text_frame.paragraphs) == 0:
                                text_frame.add_paragraph()
                            paragraph = text_frame.paragraphs[0]
                            paragraph.text = str(x_axis_label)
                            paragraph.font.size = Pt(10)
                            paragraph.font.bold = True
                            try:
                                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                            except AttributeError:
                                pass
                            # Ensure the paragraph is not empty
                            if not paragraph.text:
                                paragraph.text = str(x_axis_label)
                            logger.info('Chart %s: Set X-axis label for scatter plot: "%s"', obj.id, x_axis_label)
                        except Exception as title_error:
                            logger.error('Chart %s: Error accessing X-axis title: %s', obj.id, title_error, exc_info=True)
                    else:
                        category_axis.has_title = False
                        logger.warning('Chart %s: No X-axis label provided for scatter plot', obj.id)
                    
                    # Additional attempt: Try to set number format through plot if available
                    try:
                        plot = chart.plots[0]
                        if hasattr(plot, 'category_axis'):
                            try:
                                plot.category_axis.tick_labels.number_format = '#,##0'
                                logger.debug('Chart %s: Set X-axis number format through plot', obj.id)
                            except (AttributeError, TypeError):
                                pass
                    except (AttributeError, IndexError, TypeError):
                        pass
                except Exception as e:
                    logger.error('Chart %s: Error setting X-axis label for scatter plot: %s', obj.id, e, exc_info=True)
                
                # Y-axis (vertical, left) - for scatter plots this is value_axis
                try:
                    value_axis = chart.value_axis
                    value_axis.has_major_gridlines = True  # Horizontal gridlines
                    value_axis.has_minor_gridlines = False
                    value_axis.tick_labels.font.size = Pt(9)
                    
                    # Format Y-axis numbers as whole numbers (no scientific notation)
                    # Try multiple format approaches for scatter plots
                    try:
                        # First try the standard format
                        value_axis.tick_labels.number_format = '#,##0'
                        logger.debug('Chart %s: Set Y-axis number format to #,##0', obj.id)
                    except AttributeError:
                        try:
                            # Alternative format without thousand separator
                            value_axis.tick_labels.number_format = '0'
                            logger.debug('Chart %s: Set Y-axis number format to 0', obj.id)
                        except AttributeError:
                            try:
                                # Try General format as fallback
                                value_axis.tick_labels.number_format = 'General'
                                logger.debug('Chart %s: Set Y-axis number format to General', obj.id)
                            except AttributeError:
                                logger.warning('Chart %s: Could not set Y-axis number format', obj.id)
                    
                    if y_axis_label:
                        value_axis.has_title = True
                        # Ensure axis title exists and set text
                        try:
                            axis_title = value_axis.axis_title
                            text_frame = axis_title.text_frame
                            text_frame.clear()
                            # Ensure at least one paragraph exists
                            if len(text_frame.paragraphs) == 0:
                                text_frame.add_paragraph()
                            paragraph = text_frame.paragraphs[0]
                            paragraph.text = str(y_axis_label)
                            paragraph.font.size = Pt(10)
                            paragraph.font.bold = True
                            try:
                                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                            except AttributeError:
                                pass
                            # Ensure the paragraph is not empty
                            if not paragraph.text:
                                paragraph.text = str(y_axis_label)
                            
                            # Force update - sometimes need to access the title again
                            try:
                                # Verify the title is set
                                if value_axis.has_title and axis_title.text_frame.paragraphs[0].text:
                                    logger.info('Chart %s: Set Y-axis label for scatter plot: "%s" (verified)', obj.id, y_axis_label)
                                else:
                                    logger.warning('Chart %s: Y-axis label may not be set correctly', obj.id)
                            except Exception:
                                pass
                        except Exception as title_error:
                            logger.error('Chart %s: Error accessing Y-axis title: %s', obj.id, title_error, exc_info=True)
                    else:
                        value_axis.has_title = False
                        logger.warning('Chart %s: No Y-axis label provided for scatter plot', obj.id)
                    
                    if bool(axis_includes_zero):
                        try:
                            value_axis.crosses_at = 0
                        except AttributeError:
                            pass
                    
                    # Additional attempt: Try to set number format through plot if available
                    # Sometimes scatter plots need format set through the plot object
                    try:
                        plot = chart.plots[0]
                        # Try to access axes through plot
                        if hasattr(plot, 'value_axis'):
                            try:
                                plot.value_axis.tick_labels.number_format = '#,##0'
                                logger.debug('Chart %s: Set Y-axis number format through plot', obj.id)
                            except (AttributeError, TypeError):
                                pass
                    except (AttributeError, IndexError, TypeError):
                        pass
                    
                    # Force refresh of axis formatting - sometimes need to re-access
                    try:
                        # Re-access the axis to ensure formatting is applied
                        refreshed_axis = chart.value_axis
                        if hasattr(refreshed_axis, 'tick_labels'):
                            try:
                                refreshed_axis.tick_labels.number_format = '#,##0'
                            except AttributeError:
                                pass
                    except Exception:
                        pass
                except Exception as e:
                    logger.error('Chart %s: Error setting Y-axis label for scatter plot: %s', obj.id, e, exc_info=True)
            else:
                # Regular charts (bar, line, etc.) - category axis for X, value axis for Y
                # Category axis (X-axis) - no vertical gridlines
                category_axis = chart.category_axis
                category_axis.has_major_gridlines = False  # No vertical gridlines
                category_axis.has_minor_gridlines = False
                category_axis.tick_labels.font.size = Pt(9)
                if x_axis_label:
                    category_axis.has_title = True
                    category_axis.axis_title.text_frame.text = str(x_axis_label)
                    # Set font size and bold for axis title
                    try:
                        category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                        category_axis.axis_title.text_frame.paragraphs[0].font.bold = True
                    except (AttributeError, IndexError):
                        pass
                else:
                    category_axis.has_title = False
                
                # Value axis (Y-axis) - horizontal gridlines only
                value_axis = chart.value_axis
                value_axis.has_major_gridlines = True  # Horizontal gridlines
                value_axis.has_minor_gridlines = False
                value_axis.tick_labels.font.size = Pt(9)
                
                # Format axis numbers as whole numbers (no scientific notation)
                # Use '#,##0' format: whole numbers with thousand separators
                try:
                    value_axis.tick_labels.number_format = '#,##0'
                except AttributeError:
                    pass
                
                if y_axis_label:
                    value_axis.has_title = True
                    value_axis.axis_title.text_frame.text = str(y_axis_label)
                    # Set font size and bold for axis title
                    try:
                        value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                        value_axis.axis_title.text_frame.paragraphs[0].font.bold = True
                    except (AttributeError, IndexError):
                        pass
                else:
                    value_axis.has_title = False
                
                if bool(axis_includes_zero):
                    value_axis.crosses_at = 0
            
            # Handle dual axis: assign series to secondary axis if configured
            if use_dual_axis and series_to_axis:
                try:
                    # Get or create secondary value axis
                    # Access through plot for better compatibility
                    plot = chart.plots[0]
                    secondary_value_axis = chart.secondary_value_axis
                    secondary_value_axis.has_major_gridlines = True
                    secondary_value_axis.has_minor_gridlines = False
                    secondary_value_axis.tick_labels.font.size = Pt(9)
                    
                    # Format secondary axis numbers as whole numbers
                    try:
                        secondary_value_axis.tick_labels.number_format = '#,##0'
                    except AttributeError:
                        pass
                    
                    # Set secondary axis label - always set to make it visible
                    secondary_axis_label = str(second_y_axis_field) if second_y_axis_field else 'Secondary Axis'
                    secondary_value_axis.has_title = True
                    secondary_value_axis.axis_title.text_frame.text = secondary_axis_label
                    try:
                        secondary_value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                        # Ensure the title is visible by setting font properties
                        secondary_value_axis.axis_title.text_frame.paragraphs[0].font.bold = True
                        try:
                            secondary_value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                        except AttributeError:
                            pass
                    except (AttributeError, IndexError) as e:
                        logger.warning('Chart %s: Could not set secondary axis title font: %s', obj.id, e)
                    logger.info('Chart %s: Secondary axis title set to: %s', obj.id, secondary_axis_label)
                    
                    # Assign series to appropriate axes
                    # For dual axis, we need to set axis_group on the series
                    # This works for line, column, and bar charts
                    for series_index, series in enumerate(chart.series):
                        if series_index in series_to_axis:
                            axis_group = series_to_axis[series_index]
                            if axis_group == 'secondary':
                                # Set axis_group to 2 for secondary axis
                                # This is the standard way in python-pptx
                                try:
                                    # Direct assignment to series
                                    series.axis_group = 2
                                    logger.info('Chart %s: Series %d "%s" assigned to secondary axis', 
                                               obj.id, series_index, series.name)
                                except AttributeError:
                                    # Try through plot if direct assignment doesn't work
                                    try:
                                        if hasattr(plot, 'series') and series_index < len(plot.series):
                                            plot.series[series_index].axis_group = 2
                                            logger.info('Chart %s: Series %d "%s" assigned to secondary axis (via plot)', 
                                                       obj.id, series_index, series.name)
                                        else:
                                            # Try XML manipulation as last resort
                                            logger.warning('Chart %s: Series %d "%s" - axis_group not available, may need manual setup in PowerPoint', 
                                                          obj.id, series_index, series.name)
                                    except (AttributeError, IndexError, Exception) as e:
                                        logger.warning('Chart %s: Could not assign series %d "%s" to secondary axis: %s', 
                                                      obj.id, series_index, series.name, e)
                except AttributeError as exc:
                    logger.warning('Chart %s: Unable to configure dual axis: %s. Chart type may not support dual axis.', obj.id, exc)
                except Exception as exc:
                    logger.error('Chart %s: Error configuring dual axis: %s', obj.id, exc, exc_info=True)
            
            # Adjust plot area to position legend at the bottom border
            # The legend should stick to the bottom edge of the chart border
            try:
                plot_area = chart.plots[0].plot_area
                if chart.has_legend and chart.legend.position == XL_LEGEND_POSITION.BOTTOM:
                    # With reduced chart height and include_in_layout = True,
                    # PowerPoint will position the legend at the bottom border of the chart
                    # The plot area will automatically adjust to make room
                    pass
            except AttributeError:
                pass
        except AttributeError as exc:
            logger.debug('Chart type %s axis configuration: %s', chart_type, exc)

    for index, series in enumerate(chart.series):
        if index < len(series_colors):
            rgb = _parse_hex_color(series_colors[index])
            if rgb is not None:
                try:
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = rgb
                except AttributeError:
                    pass

    rotation = _safe_float(obj.rotation, 0)
    if rotation:
        chart_shape.rotation = rotation


def _render_atom(slide, obj: SlideExportObjectPayload, offset_x: float = 0.0, offset_y: float = 0.0, chart_objects: Optional[dict[str, SlideExportObjectPayload]] = None) -> None:
    width = _safe_float(obj.width, 0)
    height = _safe_float(obj.height, 0)
    if width <= 0 or height <= 0:
        logger.warning('Atom %s: Invalid dimensions (width=%s, height=%s)', obj.id, width, height)
        return

    props = obj.props or {}
    atom = _as_dict(props.get('atom'))
    if not atom:
        logger.warning('Atom object %s: Missing atom payload in props. Props keys: %s', obj.id, list(props.keys())[:10])
        return

    metadata = _as_dict(atom.get('metadata')) or {}
    
    # Log comprehensive debugging info
    logger.info('Atom %s: Rendering atom. Title: %s, Category: %s', 
                obj.id, atom.get('title') or atom.get('name'), atom.get('category'))
    if metadata:
        logger.info('Atom %s metadata keys (%d total): %s', obj.id, len(metadata.keys()), list(metadata.keys())[:20])
    else:
        logger.warning('Atom %s: No metadata found', obj.id)

    # REMOVED: Blue background atom container (rounded rectangle)
    # The user requested to remove the blue background container
    # No background shape is rendered - charts and content render directly on the slide

    padding = 24.0
    inner_width = max(width - (padding * 2), 0)
    if inner_width <= 0:
        return

    # Extract chart/table preview first to determine if we should hide title/category
    table_preview = _extract_table_preview(metadata)
    chart_preview = _extract_chart_preview(metadata)
    
    # Also check atom props directly for chart data
    if not chart_preview:
        atom_chart_data = atom.get('chartData') or atom.get('chart_data') or atom.get('chart')
        if atom_chart_data:
            logger.debug('Atom %s: Found chart data in atom props', obj.id)
            chart_preview = _extract_chart_preview(_as_dict(atom_chart_data) or {'data': atom_chart_data})
    
    # Also check object props for chart data
    if not chart_preview:
        obj_chart_data = props.get('chartData') or props.get('chart_data') or props.get('chart')
        if obj_chart_data:
            logger.debug('Atom %s: Found chart data in object props', obj.id)
            chart_preview = _extract_chart_preview(_as_dict(obj_chart_data) or {'data': obj_chart_data})
    
    if chart_preview:
        logger.info('Atom %s: Extracted chart preview (type: %s, categories: %d, series: %d)', 
                    obj.id, chart_preview.get('type'), 
                    len(chart_preview.get('categories', [])),
                    len(chart_preview.get('series', [])))
    else:
        logger.warning('Atom %s: No chart preview extracted from metadata, atom props, or object props. Metadata keys: %s', 
                      obj.id, list(metadata.keys())[:20] if metadata else 'none')

    # Check for post-animation images in multiple places:
    # 1. Atom metadata (from backend)
    # 2. Atom props (from frontend chart snapshots)
    # 3. Object props (fallback)
    # 4. Related chart objects (if atom has chart children)
    chart_overlay_props: dict[str, Any] = {}
    for key in (
        'postAnimationPng',
        'postAnimationSvg',
        'postAnimationImage',  # Alternative key name
        'postAnimationWidth',
        'postAnimationHeight',
        'postAnimationPixelRatio',
    ):
        # Check in order: metadata -> atom props -> object props -> related chart objects
        value = metadata.get(key)
        if value is None:
            value = atom.get(key)
        if value is None and key in props:
            value = props.get(key)
        # Check related chart objects if available
        if value is None and chart_objects:
            atom_id = atom.get('atomId') or atom.get('id')
            if atom_id:
                # Look for chart objects that might be related (e.g., same ID prefix)
                for chart_id, chart_obj in chart_objects.items():
                    chart_props = chart_obj.props or {}
                    if atom_id in chart_id or chart_id in str(atom_id):
                        value = chart_props.get(key)
                        if value:
                            logger.info('Atom %s: Found %s in related chart object %s', obj.id, key, chart_id)
                            break
        if value is not None:
            chart_overlay_props[key] = value
            logger.info('Atom %s: Found %s', obj.id, key)

    preview_image = None
    for key in ('previewImage', 'preview_image', 'image', 'thumbnail'):
        candidate = metadata.get(key)
        if _is_non_empty_str(candidate) and _is_data_url(str(candidate)):
            preview_image = str(candidate)
            break

    # Check if we have a chart - if so, hide title/category to avoid redundancy
    has_chart = bool(chart_preview or chart_overlay_props or preview_image)
    
    # Only show title/category if there's no chart (charts are self-explanatory)
    title_text = ''
    if not has_chart:
        title_text = str(atom.get('title') or atom.get('name') or '').strip()
        if not title_text:
            raw_id = atom.get('atomId') or atom.get('id')
            if _is_non_empty_str(raw_id):
                title_text = _humanise_key(str(raw_id))

    category_text = ''
    if not has_chart:
        category_text = str(atom.get('category') or '').strip()

    cursor_y = obj.y + padding
    remaining_height = height - (cursor_y - obj.y) - padding

    if title_text and remaining_height > 0:
        title_height = min(max(remaining_height * 0.3, 32.0), remaining_height)
        title_box = SlideExportObjectPayload.model_validate(
            {
                'id': f'{obj.id}::title',
                'type': 'text-box',
                'x': obj.x + padding,
                'y': cursor_y,
                'width': inner_width,
                'height': title_height,
                'props': {
                    'text': title_text,
                    'fontFamily': props.get('fontFamily') or 'Inter',
                    'fontSize': 24,
                    'bold': True,
                    'align': 'left',
                    'color': '#111827',
                },
            }
        )
        _render_text_box(slide, title_box, offset_x, offset_y)
        cursor_y += title_height + 6
        remaining_height = height - (cursor_y - obj.y) - padding

    if category_text and remaining_height > 0:
        category_height = min(max(remaining_height * 0.25, 24.0), remaining_height)
        category_box = SlideExportObjectPayload.model_validate(
            {
                'id': f'{obj.id}::category',
                'type': 'text-box',
                'x': obj.x + padding,
                'y': cursor_y,
                'width': inner_width,
                'height': category_height,
                'props': {
                    'text': category_text,
                    'fontFamily': props.get('fontFamily') or 'Inter',
                    'fontSize': 14,
                    'bold': False,
                    'align': 'left',
                    'color': '#4B5563',
                },
            }
        )
        _render_text_box(slide, category_box, offset_x, offset_y)
        cursor_y += category_height + 10
        remaining_height = height - (cursor_y - obj.y) - padding

    content_height = remaining_height
    if content_height <= 0:
        return

    summary_lines = _extract_summary_lines(metadata)

    summary_height = 0.0
    has_primary = bool(table_preview or chart_preview or preview_image)
    if summary_lines and has_primary:
        approx_line_height = 18.0
        estimated = len(summary_lines) * approx_line_height + 16.0
        summary_height = min(content_height * 0.35, max(estimated, 64.0))
        if summary_height > content_height - 40:
            summary_height = max(content_height * 0.25, 56.0)

    primary_height = content_height - summary_height
    if has_primary and primary_height <= 0:
        summary_height = 0.0
        primary_height = content_height

    content_start = cursor_y
    next_y = content_start

    if table_preview:
        table_cells = _build_table_cells(table_preview['headers'], table_preview['rows'])
        table_object = SlideExportObjectPayload.model_validate(
            {
                'id': f'{obj.id}::table',
                'type': 'table',
                'x': obj.x + padding,
                'y': content_start,
                'width': inner_width,
                'height': primary_height,
                'props': {'data': table_cells},
            }
        )
        _render_table(slide, table_object, offset_x, offset_y)
        next_y = content_start + primary_height
    elif chart_preview or chart_overlay_props or preview_image:
        # Prioritize post-animation images if available (best quality)
        if chart_overlay_props.get('postAnimationPng') or chart_overlay_props.get('postAnimationSvg'):
            logger.info('Atom %s: Rendering chart from post-animation image', obj.id)
            chart_props: dict[str, Any] = {}
            if chart_preview:
                chart_props['chartData'] = chart_preview
                chart_props['chartConfig'] = {
                    'type': chart_preview.get('type'),
                    'legendPosition': chart_preview.get('legendPosition'),
                    'showLegend': chart_preview.get('showLegend'),
                    'showValues': chart_preview.get('showValues'),
                    'axisIncludesZero': chart_preview.get('axisIncludesZero'),
                    'xAxisLabel': chart_preview.get('xAxisLabel'),
                    'yAxisLabel': chart_preview.get('yAxisLabel'),
                    'yAxis': chart_preview.get('yAxis'),
                    'secondYAxis': chart_preview.get('secondYAxis'),
                    'dualAxisMode': chart_preview.get('dualAxisMode'),
                    'seriesColors': chart_preview.get('seriesColors'),
                }
            chart_props.update(chart_overlay_props)
            
            chart_object = SlideExportObjectPayload.model_validate(
                {
                    'id': f'{obj.id}::chart',
                    'type': 'chart',
                    'x': obj.x + padding,
                    'y': content_start,
                    'width': inner_width,
                    'height': primary_height,
                    'props': chart_props,
                }
            )
            try:
                _render_chart(slide, chart_object, offset_x, offset_y)
                logger.info('Atom %s: Successfully rendered chart from post-animation image', obj.id)
            except Exception as exc:
                logger.error('Atom %s: Failed to render chart from post-animation image: %s', obj.id, exc, exc_info=True)
                # Try preview image as fallback
                if preview_image:
                    logger.info('Atom %s: Falling back to preview image', obj.id)
                    image_object = SlideExportObjectPayload.model_validate(
                        {
                            'id': f'{obj.id}::chart-fallback-image',
                            'type': 'image',
                            'x': obj.x + padding,
                            'y': content_start,
                            'width': inner_width,
                            'height': primary_height,
                            'props': {'dataUrl': preview_image},
                        }
                    )
                    _render_image(slide, image_object, offset_x, offset_y)
        elif chart_preview:
            # Render from chart data (native PowerPoint chart)
            logger.info('Atom %s: Rendering chart from chart data (no post-animation image)', obj.id)
            chart_props: dict[str, Any] = {
                'chartData': chart_preview,
                'chartConfig': {
                    'type': chart_preview.get('type'),
                    'legendPosition': chart_preview.get('legendPosition'),
                    'showLegend': chart_preview.get('showLegend'),
                    'showValues': chart_preview.get('showValues'),
                    'axisIncludesZero': chart_preview.get('axisIncludesZero'),
                    'xAxisLabel': chart_preview.get('xAxisLabel'),
                    'yAxisLabel': chart_preview.get('yAxisLabel'),
                    'yAxis': chart_preview.get('yAxis'),
                    'secondYAxis': chart_preview.get('secondYAxis'),
                    'dualAxisMode': chart_preview.get('dualAxisMode'),
                    'seriesColors': chart_preview.get('seriesColors'),
                }
            }
            chart_props.update(chart_overlay_props)
            
            chart_object = SlideExportObjectPayload.model_validate(
                {
                    'id': f'{obj.id}::chart',
                    'type': 'chart',
                    'x': obj.x + padding,
                    'y': content_start,
                    'width': inner_width,
                    'height': primary_height,
                    'props': chart_props,
                }
            )
            try:
                _render_chart(slide, chart_object, offset_x, offset_y)
                logger.info('Atom %s: Successfully rendered chart from chart data', obj.id)
            except Exception as exc:
                logger.error('Atom %s: Failed to render chart from chart data: %s', obj.id, exc, exc_info=True)
                # Try preview image as fallback
                if preview_image:
                    logger.info('Atom %s: Falling back to preview image after chart data failure', obj.id)
                    image_object = SlideExportObjectPayload.model_validate(
                        {
                            'id': f'{obj.id}::chart-fallback-image2',
                            'type': 'image',
                            'x': obj.x + padding,
                            'y': content_start,
                            'width': inner_width,
                            'height': primary_height,
                            'props': {'dataUrl': preview_image},
                        }
                    )
                    _render_image(slide, image_object, offset_x, offset_y)
        elif preview_image:
            # Last resort: use preview image
            logger.info('Atom %s: Rendering chart as preview image (no chart data or post-animation image)', obj.id)
            image_object = SlideExportObjectPayload.model_validate(
                {
                    'id': f'{obj.id}::chart-preview',
                    'type': 'image',
                    'x': obj.x + padding,
                    'y': content_start,
                    'width': inner_width,
                    'height': primary_height,
                    'props': {'dataUrl': preview_image},
                }
            )
            _render_image(slide, image_object, offset_x, offset_y)
        else:
            logger.warning('Atom %s: Has chart indication but no chart data, post-animation image, or preview image', obj.id)
        next_y = content_start + primary_height
    elif preview_image:
        image_object = SlideExportObjectPayload.model_validate(
            {
                'id': f'{obj.id}::image',
                'type': 'image',
                'x': obj.x + padding,
                'y': content_start,
                'width': inner_width,
                'height': primary_height,
                'props': {'dataUrl': preview_image},
            }
        )
        _render_image(slide, image_object, offset_x, offset_y)
        next_y = content_start + primary_height
    elif summary_lines:
        summary_height = content_height
        next_y = content_start

    if summary_lines:
        if has_primary and summary_height > 0:
            next_y += 12
        summary_text = '\n'.join(f'• {line}' for line in summary_lines)
        summary_object = SlideExportObjectPayload.model_validate(
            {
                'id': f'{obj.id}::summary',
                'type': 'text-box',
                'x': obj.x + padding,
                'y': next_y,
                'width': inner_width,
                'height': summary_height if summary_height > 0 else content_height,
                'props': {
                    'text': summary_text,
                    'fontFamily': props.get('fontFamily') or 'Inter',
                    'fontSize': 12,
                    'align': 'left',
                    'color': '#1F2937',
                },
            }
        )
        _render_text_box(slide, summary_object, offset_x, offset_y)



def _sort_objects(objects: Iterable[SlideExportObjectPayload]) -> list[SlideExportObjectPayload]:
    return sorted(objects, key=lambda item: (item.z_index if item.z_index is not None else 0))


def _render_slide_objects(
    slide,
    slide_payload: SlideExportPayload,
    offset_x: float = 0.0,
    offset_y: float = 0.0,
) -> None:
    # First pass: collect chart objects that might be related to atoms
    chart_objects_by_id: dict[str, SlideExportObjectPayload] = {}
    for obj in slide_payload.objects:
        if obj.type == 'chart':
            chart_objects_by_id[obj.id] = obj
            logger.info('Found chart object %s on slide %s', obj.id, slide_payload.id)
    
    # Second pass: render all objects
    for obj in _sort_objects(slide_payload.objects):
        try:
            if obj.type == 'text-box':
                _render_text_box(slide, obj, offset_x, offset_y)
            elif obj.type == 'image':
                # Image rendering is now graceful - it will skip failed images instead of raising
                try:
                    _render_image(slide, obj, offset_x, offset_y)
                except Exception as exc:
                    # Image errors are already logged in _render_image, just continue
                    logger.warning('Image %s failed to render, continuing with export: %s', obj.id, exc)
            elif obj.type == 'table':
                _render_table(slide, obj, offset_x, offset_y)
            elif obj.type == 'chart':
                logger.info('Rendering chart object %s on slide %s', obj.id, slide_payload.id)
                _render_chart(slide, obj, offset_x, offset_y)
            elif obj.type == 'shape':
                _render_shape(slide, obj, offset_x, offset_y)
            elif obj.type == 'atom':
                logger.info('Rendering atom object %s on slide %s', obj.id, slide_payload.id)
                # Pass chart objects map so atom can find related charts
                _render_atom(slide, obj, offset_x, offset_y, chart_objects_by_id)
            else:
                logger.debug('Skipping unsupported object type %s on slide %s', obj.type, slide_payload.id)
        except ExportGenerationError:
            raise
        except Exception as exc:  # pragma: no cover - best effort logging
            logger.exception('Failed to render %s (type: %s) on slide %s: %s', obj.id, obj.type, slide_payload.id, exc)


def _render_screenshot_background(
    slide,
    slide_payload: SlideExportPayload,
    base_width: float,
    base_height: float,
    offset_x: float,
    offset_y: float,
) -> bool:
    screenshot = slide_payload.screenshot
    if not isinstance(screenshot, SlideScreenshotPayload):
        return False

    try:
        image_bytes = _decode_data_url(screenshot.data_url)
    except ExportGenerationError as exc:
        logger.warning('Unable to decode screenshot for slide %s: %s', slide_payload.id, exc)
        return False
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.warning('Unexpected screenshot decoding error for slide %s: %s', slide_payload.id, exc)
        return False

    image_stream = io.BytesIO(image_bytes)
    try:
        picture = slide.shapes.add_picture(
            image_stream,
            _px_to_emu(offset_x),
            _px_to_emu(offset_y),
            width=_px_to_emu(base_width),
            height=_px_to_emu(base_height),
        )
    except Exception as exc:  # pragma: no cover - drawing edge cases
        logger.warning('Unable to add screenshot background for slide %s: %s', slide_payload.id, exc)
        return False

    try:
        picture.name = f"{slide_payload.id}-background"
    except Exception:  # pragma: no cover - property best effort
        pass

    try:
        lock = picture.lock
        lock.aspect_ratio = True
        lock.position = True
        lock.rotation = True
        lock.crop = True
    except AttributeError:
        pass

    return True


def _resolve_slide_dimensions(slide: SlideExportPayload) -> tuple[float, float]:
    width = _safe_float(slide.base_width, 0)
    height = _safe_float(slide.base_height, 0)

    if (width <= 0 or height <= 0) and getattr(slide, "dom_snapshot", None):
        snapshot = slide.dom_snapshot
        if isinstance(snapshot, SlideDomSnapshotPayload):
            width = max(width, _safe_float(getattr(snapshot, "width", None), 0))
            height = max(height, _safe_float(getattr(snapshot, "height", None), 0))

    if (width <= 0 or height <= 0) and slide.screenshot:
        screenshot = slide.screenshot
        width = max(width, _safe_float(getattr(screenshot, 'css_width', None), 0))
        height = max(height, _safe_float(getattr(screenshot, 'css_height', None), 0))
        width = width or _safe_float(getattr(screenshot, 'width', None), 0)
        height = height or _safe_float(getattr(screenshot, 'height', None), 0)

    if width <= 0 or height <= 0:
        raise ExportGenerationError('Slide dimensions are missing or invalid.')

    return width, height


def _build_slide_metadata(slide_payload: SlideExportPayload) -> dict[str, Any]:
    objects: list[dict[str, Any]] = []
    for obj in slide_payload.objects:
        entry: dict[str, Any] = {
            "id": obj.id,
            "type": obj.type,
            "x": obj.x,
            "y": obj.y,
            "width": obj.width,
            "height": obj.height,
            "rotation": obj.rotation,
            "zIndex": obj.z_index,
            "props": obj.props or {},
        }
        group_id = getattr(obj, "group_id", None)
        if group_id is not None:
            entry["groupId"] = group_id
        objects.append(entry)

    metadata: dict[str, Any] = {
        "id": slide_payload.id,
        "index": slide_payload.index,
        "title": slide_payload.title,
        "baseWidth": slide_payload.base_width,
        "baseHeight": slide_payload.base_height,
        "presentationSettings": slide_payload.presentation_settings,
        "objects": objects,
    }

    screenshot = slide_payload.screenshot
    if isinstance(screenshot, SlideScreenshotPayload):
        metadata["screenshot"] = {
            "width": screenshot.width,
            "height": screenshot.height,
            "cssWidth": getattr(screenshot, "css_width", None),
            "cssHeight": getattr(screenshot, "css_height", None),
            "pixelRatio": getattr(screenshot, "pixel_ratio", None),
        }

    return metadata


def _attach_slide_metadata(slide, slide_payload: SlideExportPayload) -> None:
    metadata = _build_slide_metadata(slide_payload)
    try:
        serialised = json.dumps(metadata, ensure_ascii=False, separators=(",", ":"))
    except (TypeError, ValueError) as exc:
        logger.warning('Unable to serialise metadata for slide %s: %s', slide_payload.id, exc)
        return

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    existing_text = [paragraph.text for paragraph in text_frame.paragraphs if paragraph.text.strip()]

    text_frame.clear()

    for paragraph_text in existing_text:
        paragraph = text_frame.add_paragraph()
        paragraph.text = paragraph_text

    marker = text_frame.add_paragraph()
    marker.text = METADATA_MARKER
    marker.level = 0
    try:
        marker.font.size = Pt(6)
    except AttributeError:
        pass

    data_paragraph = text_frame.add_paragraph()
    data_paragraph.text = serialised
    data_paragraph.level = 0
    try:
        data_paragraph.font.size = Pt(6)
    except AttributeError:
        pass


def _prepare_render_slide(slide: SlideExportPayload) -> dict:
    snapshot = slide.dom_snapshot
    if not isinstance(snapshot, SlideDomSnapshotPayload):
        raise ExportGenerationError(
            f'Slide {slide.id} is missing a DOM snapshot required for server-side rendering.'
        )

    width = _safe_float(getattr(snapshot, "width", None), 0) or _safe_float(slide.base_width, 0)
    height = _safe_float(getattr(snapshot, "height", None), 0) or _safe_float(slide.base_height, 0)

    if width <= 0 or height <= 0:
        raise ExportGenerationError(
            f'Slide {slide.id} does not include valid dimensions for rendering.'
        )

    payload: dict[str, object] = {
        "id": slide.id,
        "html": snapshot.html,
        "width": width,
        "height": height,
    }

    pixel_ratio = _safe_float(getattr(snapshot, "pixel_ratio", None), 0)
    if pixel_ratio > 0:
        payload["pixelRatio"] = pixel_ratio

    return payload


def _request_slide_screenshots(
    slides: Sequence[SlideExportPayload],
    styles: DocumentStylesPayload,
    *,
    strict: bool = True,
) -> dict[str, dict]:
    if not slides:
        return {}

    render_slides = []
    pixel_ratios: list[float] = []
    for slide in slides:
        render_payload = _prepare_render_slide(slide)
        ratio = _safe_float(render_payload.get("pixelRatio"), 0)
        if ratio > 0:
            pixel_ratios.append(ratio)
        render_slides.append(render_payload)

    try:
        inputs = build_inputs(render_slides)
        rendered = render_slide_batch(
            inputs,
            styles,
            pixel_ratio=max(pixel_ratios) if pixel_ratios else None,
        )
    except ExhibitionRendererError as exc:
        message = f"Server-side renderer failed to capture slides: {exc}"
        if strict:
            raise ExportGenerationError(message) from exc
        logger.warning(message)
        return {}

    return {entry.id: entry.as_payload() for entry in rendered}


def _attempt_server_screenshots(
    payload: ExhibitionExportRequest, slides: Sequence[SlideExportPayload]
) -> None:
    """Attempt server-side screenshots, but skip if client-side screenshots are already available.
    
    NEW APPROACH: Frontend now captures visible exhibition slides directly as images.
    Server-side screenshots are only used as fallback when client-side captures are missing.
    """
    if not slides:
        return

    # NEW APPROACH: Skip server-side rendering if slides already have client-side screenshots
    # These are captured directly from visible exhibition, so they're already complete
    slides_with_screenshots = [
        slide for slide in slides 
        if slide.screenshot and isinstance(slide.screenshot, SlideScreenshotPayload) and slide.screenshot.data_url
    ]
    
    if len(slides_with_screenshots) == len(slides):
        logger.info('All slides have client-side screenshots from visible exhibition, skipping server-side rendering')
        return

    # Only attempt server-side for slides missing screenshots
    styles = payload.document_styles
    if not isinstance(styles, DocumentStylesPayload):
        logger.debug('Document styles missing, cannot perform server-side rendering')
        return

    candidates = [
        slide
        for slide in slides
        if not (slide.screenshot and isinstance(slide.screenshot, SlideScreenshotPayload) and slide.screenshot.data_url)
        and isinstance(slide.dom_snapshot, SlideDomSnapshotPayload) 
        and slide.dom_snapshot.html
    ]
    
    if not candidates:
        logger.debug('No slides need server-side rendering')
        return

    logger.info('Attempting server-side rendering for %d slide(s) missing client-side screenshots', len(candidates))

    try:
        screenshots = _request_slide_screenshots(candidates, styles, strict=False)
    except ExportGenerationError as exc:  # pragma: no cover - best effort logging
        logger.warning('Falling back to client slide captures: %s', exc)
        return

    for slide in candidates:
        data = screenshots.get(slide.id)
        if not isinstance(data, dict):
            continue
        try:
            slide.screenshot = SlideScreenshotPayload.model_validate(data)
            logger.info('Slide %s: Server-side screenshot captured', slide.id)
        except Exception as exc:  # pragma: no cover - validation edge cases
            logger.warning('Skipping invalid renderer screenshot for slide %s: %s', slide.id, exc)


def _ensure_slide_screenshots(
    payload: ExhibitionExportRequest, slides: Sequence[SlideExportPayload]
) -> None:
    missing = [slide for slide in slides if not slide.screenshot or not slide.screenshot.data_url]
    if not missing:
        return

    styles = payload.document_styles
    if not isinstance(styles, DocumentStylesPayload):
        raise ExportGenerationError(
            'Document styles are required to render slide screenshots on the server.'
        )

    screenshots = _request_slide_screenshots(missing, styles)

    for slide in missing:
        data = screenshots.get(slide.id)
        if not isinstance(data, dict):
            raise ExportGenerationError(f'Unable to render screenshot for slide {slide.id}.')
        slide.screenshot = SlideScreenshotPayload.model_validate(data)


def render_slide_screenshots(payload: ExhibitionExportRequest) -> list[dict[str, object]]:
    """Render slide screenshots using the server-side rendering service."""

    if not payload.slides:
        raise ExportGenerationError('No slides provided for export.')

    ordered_slides = sorted(payload.slides, key=lambda slide: slide.index)

    _attempt_server_screenshots(payload, ordered_slides)
    _ensure_slide_screenshots(payload, ordered_slides)

    rendered: list[SlideScreenshotResponse] = []
    for slide in ordered_slides:
        _apply_post_animation_overlays(slide)
        screenshot = slide.screenshot
        if not isinstance(screenshot, SlideScreenshotPayload):
            raise ExportGenerationError(
                f'Unable to render screenshot for slide {slide.id}.'
            )

        entry = SlideScreenshotResponse.model_validate(
            {
                "id": slide.id,
                "index": slide.index,
                **screenshot.model_dump(by_alias=True),
            }
        )
        rendered.append(entry)

    return [entry.model_dump(by_alias=True) for entry in rendered]


def build_pptx_bytes_animated(payload: ExhibitionExportRequest) -> bytes:
    """Build PPTX export with animation preservation.
    
    This function uses object-based rendering instead of static screenshots,
    which preserves animations and dynamic chart interactions in PowerPoint.
    
    IMPORTANT: This is a separate pipeline from PDF/JPG exports.
    - Does NOT use Chromium screenshots
    - Does NOT call screenshot capture functions
    - Renders objects directly to preserve animations
    - Data is loaded through metadata attachment
    """
    if not payload.slides:
        raise ExportGenerationError('No slides provided for export.')

    ordered_slides = sorted(payload.slides, key=lambda slide: slide.index)
    logger.info('Starting PPTX export: %d slide(s), title: %s', len(ordered_slides), payload.title or 'Untitled')
    
    # NOTE: We do NOT call screenshot functions here - this preserves animations
    # _attempt_server_screenshots() and _ensure_slide_screenshots() are NOT called
    # This is intentional - we want object-based rendering, not static screenshots
    
    # Count image objects for logging
    total_image_objects = sum(
        len([obj for obj in slide.objects if obj.type == 'image'])
        for slide in ordered_slides
    )
    logger.info('PPTX export: Found %d image object(s) across %d slide(s)', total_image_objects, len(ordered_slides))
    
    dimensions = [_resolve_slide_dimensions(slide) for slide in ordered_slides]
    max_width = max(width for width, _ in dimensions)
    max_height = max(height for _, height in dimensions)

    presentation = Presentation()
    presentation.slide_width = _px_to_emu(max_width)
    presentation.slide_height = _px_to_emu(max_height)

    title = (payload.title or 'Exhibition Presentation').strip() or 'Exhibition Presentation'
    presentation.core_properties.title = title
    presentation.core_properties.subject = 'Exhibition export'
    presentation.core_properties.author = 'Trinity Exhibition'

    for slide_payload, (base_width, base_height) in zip(ordered_slides, dimensions):
        logger.debug('PPTX: Processing slide %d/%d: %s', 
                    slide_payload.index + 1, len(ordered_slides), slide_payload.id)
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        offset_x = max((max_width - base_width) / 2, 0.0)
        offset_y = max((max_height - base_height) / 2, 0.0)
        
        # Count objects for this slide
        slide_image_objects = [obj for obj in slide_payload.objects if obj.type == 'image']
        logger.debug('PPTX: Slide %s has %d image object(s) to render', slide_payload.id, len(slide_image_objects))
        
        # Object-based rendering (preserves animations)
        # This renders actual chart objects, shapes, text, etc. instead of static screenshots
        try:
            _render_layout_overlay(slide, slide_payload, base_width, base_height, offset_x, offset_y)
            _render_slide_objects(slide, slide_payload, offset_x, offset_y)
            
            # Attach metadata for data loading (charts can access their data from metadata)
            _attach_slide_metadata(slide, slide_payload)
            logger.debug('PPTX: Successfully processed slide %s', slide_payload.id)
        except Exception as exc:
            logger.error('PPTX: Failed to render slide %s: %s', slide_payload.id, exc, exc_info=True)
            # Continue with other slides instead of failing entire export
            logger.warning('PPTX: Continuing with remaining slides despite error on slide %s', slide_payload.id)

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output.getvalue()


# Backward compatibility alias - tests and legacy code may still reference this
build_pptx_bytes = build_pptx_bytes_animated


def build_pdf_bytes(payload: ExhibitionExportRequest) -> bytes:
    """Build PDF export from slides with screenshots.
    
    PDF export uses Chromium screenshots (same as JPG export), so all user images
    are already captured in the screenshots. We skip processing individual image
    objects and only overlay post-animation chart images if needed.
    """
    if not payload.slides:
        raise ExportGenerationError('No slides provided for export.')

    ordered_slides = sorted(payload.slides, key=lambda slide: slide.index)
    logger.info('Starting PDF export: %d slide(s), title: %s', len(ordered_slides), payload.title or 'Untitled')

    # Ensure screenshots are available (same as JPG export)
    logger.debug('Ensuring screenshots for PDF export')
    _attempt_server_screenshots(payload, ordered_slides)
    _ensure_slide_screenshots(payload, ordered_slides)
    
    # Verify all slides have screenshots
    slides_without_screenshots = [
        slide.id for slide in ordered_slides 
        if not (slide.screenshot and isinstance(slide.screenshot, SlideScreenshotPayload) and slide.screenshot.data_url)
    ]
    if slides_without_screenshots:
        logger.warning('PDF export: %d slide(s) missing screenshots: %s', 
                     len(slides_without_screenshots), slides_without_screenshots)

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.setTitle(payload.title or 'Exhibition Presentation')

    for index, slide in enumerate(ordered_slides):
        logger.debug('Processing slide %d/%d for PDF: %s', index + 1, len(ordered_slides), slide.id)
        
        width, height = _resolve_slide_dimensions(slide)
        page_width = _px_to_pt(width)
        page_height = _px_to_pt(height) 
        pdf.setPageSize((page_width, page_height))
        
        # Count objects for logging
        image_objects = [obj for obj in slide.objects if obj.type == 'image']
        atom_objects = [obj for obj in slide.objects if obj.type == 'atom']
        other_objects = [obj for obj in slide.objects if obj.type not in ('image', 'atom')]
        logger.debug('Slide %s: %d image objects, %d atom objects, %d other objects (all will be in screenshot)', 
                    slide.id, len(image_objects), len(atom_objects), len(other_objects))

        screenshot = slide.screenshot
        if not screenshot or not isinstance(screenshot.data_url, str):
            raise ExportGenerationError('Every slide must include a screenshot for PDF export.')

        # Decode and load screenshot with error handling
        try:
            logger.debug('Decoding screenshot for slide %s (PDF export)', slide.id)
            screenshot_bytes = _decode_data_url(screenshot.data_url)
            if not screenshot_bytes:
                raise ExportGenerationError(f'Screenshot for slide {slide.id} is empty after decoding.')
            
            logger.debug('Screenshot decoded: %d bytes for slide %s', len(screenshot_bytes), slide.id)
            image_stream = io.BytesIO(screenshot_bytes)
            image = ImageReader(image_stream)
            logger.debug('ImageReader created successfully for slide %s', slide.id)
        except ExportGenerationError:
            # Re-raise export errors
            raise
        except Exception as exc:
            # Wrap unexpected errors
            logger.error('Failed to decode/load screenshot for slide %s: %s', slide.id, exc, exc_info=True)
            raise ExportGenerationError(f'Failed to process screenshot for slide {slide.id}: {exc}') from exc

        css_width = _safe_float(getattr(screenshot, 'css_width', None), 0)
        css_height = _safe_float(getattr(screenshot, 'css_height', None), 0)
        pixel_ratio = _safe_float(getattr(screenshot, 'pixel_ratio', None), 0) or 1.0
        image_width = _safe_float(getattr(screenshot, 'width', None), 0)
        image_height = _safe_float(getattr(screenshot, 'height', None), 0)

        if css_width <= 0 and image_width > 0 and pixel_ratio > 0:
            css_width = image_width / pixel_ratio
        if css_height <= 0 and image_height > 0 and pixel_ratio > 0:
            css_height = image_height / pixel_ratio

        if css_width <= 0:
            css_width = width
        if css_height <= 0:
            css_height = height

        aspect_ratio = css_height / css_width if css_width > 0 else 1.0
        if aspect_ratio <= 0:
            aspect_ratio = height / width if width > 0 else 1.0

        draw_width = page_width
        draw_height = draw_width * aspect_ratio

        if draw_height > page_height and draw_height > 0:
            scale = page_height / draw_height
            draw_height = page_height
            draw_width = draw_width * scale

        offset_x = (page_width - draw_width) / 2 if draw_width < page_width else 0.0
        offset_y = (page_height - draw_height) / 2 if draw_height < page_height else 0.0

        # Draw the main screenshot (this contains all slide content including user images)
        try:
            pdf.drawImage(
                image,
                offset_x,
                offset_y,
                width=draw_width,
                height=draw_height,
                preserveAspectRatio=False,
                mask='auto',
            )
            logger.debug('Successfully drew screenshot for slide %s at (%.1f, %.1f), size: %.1f x %.1f', 
                       slide.id, offset_x, offset_y, draw_width, draw_height)
        except Exception as exc:
            logger.error('Failed to draw screenshot for slide %s: %s', slide.id, exc, exc_info=True)
            raise ExportGenerationError(f'Failed to draw screenshot for slide {slide.id}: {exc}') from exc

        scale_x = draw_width / page_width if page_width > 0 else 1.0
        scale_y = draw_height / page_height if page_height > 0 else 1.0

        # Render individual objects on top of screenshot
        # Note: Regular image objects are already captured in the screenshot, so we skip them
        # Only render post-animation images (charts) that need to be overlaid
        for obj in slide.objects:
            try:
                props = obj.props or {}
                if not isinstance(props, dict):
                    continue

                # CRITICAL: Skip ALL regular image objects - they're already in the screenshot
                # This includes user-uploaded images, preview images, and any image type objects
                if obj.type == 'image':
                    logger.debug('Skipping image object %s for PDF - already captured in screenshot', obj.id)
                    continue

                # Also skip atoms that might contain images - they're in the screenshot too
                # Atoms are complex objects that may have images, but for PDF we rely on screenshots
                if obj.type == 'atom':
                    logger.debug('Skipping atom object %s for PDF - content already in screenshot', obj.id)
                    continue

                # Only process post-animation images for charts (not regular images)
                # These are special overlay images that need to be added on top of screenshots
                image_bytes: Optional[bytes] = None
                
                try:
                    # Handle post-animation images (charts, etc.)
                    # This should only return images for chart objects, not regular user images
                    image_bytes = _resolve_post_animation_image(props)
                except Exception as exc:
                    logger.warning('Error resolving post-animation image for object %s (type: %s): %s. Skipping.', 
                                 obj.id, obj.type, exc)
                    continue
                
                if not image_bytes:
                    # No post-animation image to overlay - this is normal for most objects
                    continue

                # Log what we're about to process
                logger.debug('Processing post-animation overlay for object %s (type: %s): %d bytes', 
                           obj.id, obj.type, len(image_bytes))

                obj_width = _safe_float(obj.width, 0)
                obj_height = _safe_float(obj.height, 0)
                if obj_width <= 0 or obj_height <= 0:
                    logger.debug('Skipping object %s: invalid dimensions (%.1f x %.1f)', obj.id, obj_width, obj_height)
                    continue

                chart_x_pt = _px_to_pt(_safe_float(obj.x, 0))
                chart_y_pt = _px_to_pt(_safe_float(obj.y, 0))
                chart_width_pt = _px_to_pt(obj_width)
                chart_height_pt = _px_to_pt(obj_height)
                chart_bottom_pt = page_height - (chart_y_pt + chart_height_pt)

                draw_x = offset_x + chart_x_pt * scale_x
                draw_y = offset_y + chart_bottom_pt * scale_y
                draw_w = chart_width_pt * scale_x
                draw_h = chart_height_pt * scale_y

                # Validate dimensions before drawing
                if draw_w <= 0 or draw_h <= 0:
                    logger.warning('Skipping object %s: calculated draw dimensions invalid (%.1f x %.1f)', 
                                 obj.id, draw_w, draw_h)
                    continue

                try:
                    image_reader = ImageReader(io.BytesIO(image_bytes))
                    pdf.drawImage(
                        image_reader,
                        draw_x,
                        draw_y,
                        width=draw_w,
                        height=draw_h,
                        preserveAspectRatio=False,
                        mask='auto',
                    )
                    logger.debug('Successfully drew post-animation overlay for object %s at (%.1f, %.1f)', 
                               obj.id, draw_x, draw_y)
                except Exception as exc:
                    logger.warning('Failed to draw post-animation overlay for object %s (type: %s) on PDF: %s. Skipping.', 
                                 obj.id, obj.type, exc)
                    continue
                    
            except Exception as exc:
                # Log but continue with other objects - never fail entire export due to one object
                logger.warning('Error processing object %s (type: %s) for PDF: %s. Skipping.', obj.id, obj.type, exc)
                continue
        if index < len(ordered_slides) - 1:
            pdf.showPage()

    pdf.save()
    buffer.seek(0)
    return buffer.getvalue()


def build_export_filename(title: Optional[str], extension: str) -> str:
    base = (title or 'exhibition-export').strip().lower()
    base = re.sub(r"[^a-z0-9._-]+", '-', base)
    base = base.strip('-')[:120] or 'exhibition-export'
    return f"{base}.{extension}"
