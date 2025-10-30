from __future__ import annotations

import base64
import io
from typing import Dict, Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from .schemas import (
    BackgroundSpec,
    ChartObject,
    ExportRequest,
    ForeignObject,
    ImageObject,
    OverlaySpec,
    ShapeObject,
    SlideExportData,
    SlideScreenshot,
    TableCell,
    TableObject,
    TextObject,
)

CANVAS_WIDTH = 960
CANVAS_HEIGHT = 520
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5
DEFAULT_FONT = 'Arial'

SHAPE_TYPE_MAP: Dict[str, MSO_SHAPE] = {
    'rectangle': MSO_SHAPE.RECTANGLE,
    'rounded-rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
    'ellipse': MSO_SHAPE.OVAL,
    'circle': MSO_SHAPE.OVAL,
    'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    'diamond': MSO_SHAPE.DIAMOND,
    'pentagon': MSO_SHAPE.PENTAGON,
    'hexagon': MSO_SHAPE.HEXAGON,
    'octagon': MSO_SHAPE.OCTAGON,
    'star': MSO_SHAPE.STAR_5_POINT,
    'burst': MSO_SHAPE.EXPLOSION_1,
}

LINE_STYLE_MAP: Dict[str, MSO_LINE_DASH_STYLE] = {
    'solid': MSO_LINE_DASH_STYLE.SOLID,
    'dashed': MSO_LINE_DASH_STYLE.DASH,
    'dotted': MSO_LINE_DASH_STYLE.DOT,
}

ALIGNMENT_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
}


def _px_to_inches(value: float, axis: str) -> float:
    if axis in {'x', 'w'}:
        return (value / CANVAS_WIDTH) * SLIDE_WIDTH_IN
    return (value / CANVAS_HEIGHT) * SLIDE_HEIGHT_IN


def _px_to_points(value: float) -> float:
    return (value / 96.0) * 72.0


def _decode_base64(data: Optional[str]) -> Optional[bytes]:
    if not data:
        return None
    payload = data
    if payload.startswith('data:'):
        payload = payload.split(',', 1)[-1]
    try:
        return base64.b64decode(payload)
    except (ValueError, TypeError):
        return None


def _color_from_hex(value: Optional[str]) -> Optional[RGBColor]:
    if not value:
        return None
    hex_value = value.strip()
    if not hex_value:
        return None
    if hex_value.startswith('#'):
        hex_value = hex_value[1:]
    if len(hex_value) == 3:
        hex_value = ''.join(char * 2 for char in hex_value)
    if len(hex_value) != 6:
        return None
    try:
        return RGBColor.from_string(hex_value.upper())
    except ValueError:
        return None


def _apply_background(pres: Presentation, slide, background: BackgroundSpec) -> None:
    if background.type == 'solid' and background.color:
        fill = slide.background.fill
        fill.solid()
        color = _color_from_hex(background.color)
        if color:
            fill.fore_color.rgb = color
        return

    if background.type == 'gradient' and background.gradient and background.gradient.colors:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            pres.slide_width,
            pres.slide_height,
        )
        shape.line.fill.background()
        fill = shape.fill
        fill.gradient()
        stops = fill.gradient_stops
        start_color = _color_from_hex(background.gradient.colors[0])
        end_color = _color_from_hex(background.gradient.colors[-1])
        if start_color:
            stops[0].color.rgb = start_color
            stops[0].position = 0
        if end_color:
            stops[-1].color.rgb = end_color
            stops[-1].position = 1
        fill.gradient_angle = background.gradient.angle
        return

    if background.type == 'image':
        data = _decode_base64(background.image_data)
        if data:
            slide.shapes.add_picture(
                io.BytesIO(data),
                Inches(0),
                Inches(0),
                width=pres.slide_width,
                height=pres.slide_height,
            )
            return

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _apply_overlay(slide, overlay: OverlaySpec) -> None:
    left = Inches(_px_to_inches(overlay.x, 'x'))
    top = Inches(_px_to_inches(overlay.y, 'y'))
    width = Inches(_px_to_inches(overlay.width, 'w'))
    height = Inches(_px_to_inches(overlay.height, 'h'))

    if overlay.type == 'image':
        data = _decode_base64(overlay.image_data)
        if data:
            slide.shapes.add_picture(io.BytesIO(data), left, top, width=width, height=height)
        return

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()

    if overlay.type == 'color' and overlay.color:
        color = _color_from_hex(overlay.color)
        if color:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = color
        return

    if overlay.type == 'gradient' and overlay.gradient and overlay.gradient.colors:
        fill = shape.fill
        fill.gradient()
        stops = fill.gradient_stops
        start_color = _color_from_hex(overlay.gradient.colors[0])
        end_color = _color_from_hex(overlay.gradient.colors[-1])
        if start_color:
            stops[0].color.rgb = start_color
            stops[0].position = 0
        if end_color:
            stops[-1].color.rgb = end_color
            stops[-1].position = 1
        fill.gradient_angle = overlay.gradient.angle
        return

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _add_text(slide, obj: TextObject) -> None:
    text = obj.text or ''
    if not text.strip():
        return

    textbox = slide.shapes.add_textbox(
        Inches(_px_to_inches(obj.x, 'x')),
        Inches(_px_to_inches(obj.y, 'y')),
        Inches(_px_to_inches(obj.width, 'w')),
        Inches(_px_to_inches(obj.height, 'h')),
    )
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = obj.font_family or DEFAULT_FONT
    font.size = Pt(max(obj.font_size, 1))
    font.bold = obj.bold
    font.italic = obj.italic
    font.underline = obj.underline
    color = _color_from_hex(obj.color)
    if color:
        font.color.rgb = color
    paragraph.alignment = ALIGNMENT_MAP.get(obj.align, PP_ALIGN.LEFT)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    textbox.rotation = obj.rotation


def _add_image(slide, obj: ImageObject) -> None:
    data = _decode_base64(obj.data) or _decode_base64(obj.src)
    if not data:
        return
    slide.shapes.add_picture(
        io.BytesIO(data),
        Inches(_px_to_inches(obj.x, 'x')),
        Inches(_px_to_inches(obj.y, 'y')),
        width=Inches(_px_to_inches(obj.width, 'w')),
        height=Inches(_px_to_inches(obj.height, 'h')),
    )


def _add_shape(slide, obj: ShapeObject) -> None:
    shape_type = SHAPE_TYPE_MAP.get(obj.shape_id, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(_px_to_inches(obj.x, 'x')),
        Inches(_px_to_inches(obj.y, 'y')),
        Inches(_px_to_inches(obj.width, 'w')),
        Inches(_px_to_inches(obj.height, 'h')),
    )

    if obj.fill and obj.fill != 'transparent':
        fill = shape.fill
        fill.solid()
        color = _color_from_hex(obj.fill)
        if color:
            fill.fore_color.rgb = color
        transparency = max(0.0, min(1.0, 1.0 - obj.opacity))
        fill.transparency = transparency
    else:
        shape.fill.background()

    if obj.stroke and obj.stroke != 'transparent':
        color = _color_from_hex(obj.stroke)
        if color:
            line = shape.line
            line.fill.solid()
            line.color.rgb = color
            line.width = Pt(max(_px_to_points(obj.stroke_width), 0.25))
            dash = LINE_STYLE_MAP.get(obj.stroke_style)
            if dash:
                line.dash_style = dash
    else:
        shape.line.fill.background()

    shape.rotation = obj.rotation


def _merge_table_cell(table, row: int, col: int, row_span: int, col_span: int):
    max_row_index = len(table.rows) - 1
    max_col_index = len(table.columns) - 1
    end_row = min(max_row_index, row + max(row_span - 1, 0))
    end_col = min(max_col_index, col + max(col_span - 1, 0))
    return table.cell(row, col).merge(table.cell(end_row, end_col))


def _apply_cell_text(cell, data: TableCell) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = ALIGNMENT_MAP.get(data.formatting.align, PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = data.content or ''
    font = run.font
    font.name = data.formatting.font_family or DEFAULT_FONT
    font.size = Pt(max(data.formatting.font_size, 1))
    font.bold = data.formatting.bold
    font.italic = data.formatting.italic
    font.underline = data.formatting.underline
    color = _color_from_hex(data.formatting.color)
    if color:
        font.color.rgb = color


def _add_table(slide, obj: TableObject) -> None:
    rows = len(obj.data) if obj.data else 1
    cols = len(obj.data[0]) if rows and obj.data[0] else 1
    rows = max(rows, 1)
    cols = max(cols, 1)

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(_px_to_inches(obj.x, 'x')),
        Inches(_px_to_inches(obj.y, 'y')),
        Inches(_px_to_inches(obj.width, 'w')),
        Inches(_px_to_inches(obj.height, 'h')),
    )
    table = table_shape.table
    skipped: set[tuple[int, int]] = set()

    for row_index in range(rows):
        row_cells = obj.data[row_index] if row_index < len(obj.data) else []
        for col_index in range(cols):
            if (row_index, col_index) in skipped:
                continue
            cell_data = row_cells[col_index] if col_index < len(row_cells) else TableCell()
            cell = table.cell(row_index, col_index)
            row_span = cell_data.row_span or 1
            col_span = cell_data.col_span or 1
            if row_span > 1 or col_span > 1:
                cell = _merge_table_cell(table, row_index, col_index, row_span, col_span)
                for r in range(row_index, row_index + row_span):
                    for c in range(col_index, col_index + col_span):
                        if r == row_index and c == col_index:
                            continue
                        skipped.add((r, c))
            _apply_cell_text(cell, cell_data)


def _add_placeholder(slide, obj, message: str) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(_px_to_inches(obj.x, 'x')),
        Inches(_px_to_inches(obj.y, 'y')),
        Inches(_px_to_inches(obj.width, 'w')),
        Inches(_px_to_inches(obj.height, 'h')),
    )
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(240, 240, 240)
    textbox.line.fill.background()
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = message
    font = run.font
    font.name = DEFAULT_FONT
    font.size = Pt(14)
    font.color.rgb = RGBColor(90, 90, 90)


def build_pptx_document(payload: ExportRequest) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    presentation.core_properties.title = payload.title
    presentation.core_properties.author = 'Exhibition Mode'
    presentation.core_properties.subject = 'Exported Presentation'

    slides: Iterable[SlideExportData] = payload.slides

    for slide_data in slides:
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)
        _apply_background(presentation, slide, slide_data.background)
        if slide_data.overlay:
            _apply_overlay(slide, slide_data.overlay)

        for obj in sorted(slide_data.objects, key=lambda item: getattr(item, 'z_index', 0)):
            if isinstance(obj, TextObject):
                _add_text(slide, obj)
            elif isinstance(obj, ImageObject):
                _add_image(slide, obj)
            elif isinstance(obj, ShapeObject):
                _add_shape(slide, obj)
            elif isinstance(obj, TableObject):
                _add_table(slide, obj)
            elif isinstance(obj, ChartObject):
                _add_placeholder(slide, obj, 'Chart export is not yet supported')
            elif isinstance(obj, ForeignObject):
                _add_placeholder(slide, obj, obj.object_type or 'Embedded content')

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def build_pdf_document(payload: ExportRequest) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.setTitle(payload.title)

    screenshots: Iterable[SlideScreenshot] = payload.screenshots or []

    for screenshot in screenshots:
        data = _decode_base64(screenshot.data)
        if not data:
            pdf.showPage()
            continue
        width_px = max(screenshot.width, 1)
        height_px = max(screenshot.height, 1)
        width_in = width_px / 96.0
        height_in = height_px / 96.0
        page_width = width_in * 72.0
        page_height = height_in * 72.0
        pdf.setPageSize((page_width, page_height))
        image = ImageReader(io.BytesIO(data))
        pdf.drawImage(image, 0, 0, width=page_width, height=page_height, preserveAspectRatio=True, anchor='sw')
        pdf.showPage()

    if not payload.screenshots:
        pdf.showPage()

    pdf.save()
    buffer.seek(0)
    return buffer.getvalue()
